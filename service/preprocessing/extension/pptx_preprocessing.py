"""
PPTX 전처리 모듈
PPTX 파일을 PPT로 변환 후 텍스트와 표로 추출하는 기능 제공
"""
from __future__ import annotations
import logging
from pathlib import Path

logger = logging.getLogger(__name__)


def _clean_text(s: str | None) -> str:
    """텍스트 정규화 (공통 유틸리티에서 import)"""
    from service.preprocessing.extension.utils import clean_text as clean_text_util
    return clean_text_util(s)


def _convert_pptx_to_ppt(src: Path) -> Path | None:
    """PPTX 파일을 PPT로 변환 (LibreOffice 사용)"""
    try:
        import subprocess
        import shutil
        import time
        
        # LibreOffice 명령어 확인
        libreoffice_cmd = shutil.which("libreoffice")
        if not libreoffice_cmd:
            logger.warning(f"[PPTX->PPT] LibreOffice가 설치되지 않았습니다. PPTX 파일 변환을 건너뜁니다: {src.name}")
            return None
        
        outdir = src.parent
        logger.info(f"[PPTX->PPT] 변환 시작: {src.name} -> {outdir}")
        
        # stdout과 stderr를 캡처하여 디버깅 정보 확인
        result = subprocess.run(
            [
                libreoffice_cmd, "--headless", "--convert-to", "ppt",
                "--outdir", str(outdir), str(src)
            ],
            check=False,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            timeout=120,
        )
        
        # stdout과 stderr 로그 출력
        if result.stdout:
            stdout_msg = result.stdout.decode("utf-8", errors="ignore")
            if stdout_msg.strip():
                logger.debug(f"[PPTX->PPT] stdout: {stdout_msg[:500]}")
        
        if result.returncode != 0:
            stderr_msg = result.stderr.decode("utf-8", errors="ignore") if result.stderr else ""
            logger.error(
                f"[PPTX->PPT] 변환 실패 (returncode={result.returncode}): {src.name}. "
                f"에러: {stderr_msg[:500]}"
            )
            return None
        else:
            logger.info(f"[PPTX->PPT] LibreOffice 변환 명령 완료 (returncode=0): {src.name}")
        
        # 변환된 파일 찾기 (LibreOffice는 원본 파일명을 유지하되 확장자만 변경)
        cand = src.with_suffix(".ppt")
        
        # 파일이 생성될 때까지 잠시 대기 (비동기 처리 대비)
        max_wait = 10  # 최대 10초 대기
        wait_interval = 0.5  # 0.5초마다 확인
        waited = 0
        logger.debug(f"[PPTX->PPT] 변환된 파일 대기 중: {cand}")
        while not cand.exists() and waited < max_wait:
            time.sleep(wait_interval)
            waited += wait_interval
            if waited % 2 == 0:  # 2초마다 로그
                logger.debug(f"[PPTX->PPT] 파일 대기 중... ({waited:.1f}초)")
        
        if cand.exists():
            file_size = cand.stat().st_size
            logger.info(f"[PPTX->PPT] 변환 성공: {src.name} -> {cand.name} (크기: {file_size} bytes)")
            return cand
        else:
            # 다른 가능한 경로 확인 (대소문자 차이 등)
            parent = src.parent
            stem = src.stem
            possible_names = [
                f"{stem}.ppt",
                f"{stem}.PPT",
                f"{stem}.Ppt",
            ]
            logger.debug(f"[PPTX->PPT] 대체 경로 검색 중: {possible_names}")
            for name in possible_names:
                alt_path = parent / name
                if alt_path.exists():
                    file_size = alt_path.stat().st_size
                    logger.info(f"[PPTX->PPT] 변환 성공 (대체 경로): {src.name} -> {alt_path.name} (크기: {file_size} bytes)")
                    return alt_path
            
            # 디렉토리 전체 내용 확인
            all_files = list(parent.glob(f'{stem}.*'))
            logger.error(f"[PPTX->PPT] 변환된 파일을 찾을 수 없습니다. 예상 경로: {cand}")
            logger.error(f"[PPTX->PPT] 디렉토리 내용: {[f.name for f in all_files]}")
            logger.error(f"[PPTX->PPT] 원본 파일 존재: {src.exists()}, 경로: {src}")
            return None
            
    except subprocess.TimeoutExpired:
        logger.error(f"[PPTX->PPT] 변환 타임아웃: {src.name} (120초 초과)")
        return None
    except FileNotFoundError:
        logger.error(f"[PPTX->PPT] LibreOffice 명령어를 찾을 수 없습니다. 설치가 필요합니다.")
        return None
    except Exception as e:
        logger.exception(f"[PPTX->PPT] 변환 중 예외 발생: {src.name}, 오류: {e}")
        return None


def extract_pptx(fp: Path) -> tuple[str, list[dict]]:
    """PPTX 파일 추출 (PPT로 변환 후 처리)
    PPTX 파일이 들어오면 자동으로 PPT로 변환 후, PPT를 PPTX로 다시 변환하여 파싱합니다.
    
    Returns:
        tuple[str, list[dict]]: (본문 텍스트, 표 리스트)
    """
    # PPTX 파일인 경우 PPT로 변환
    converted_ppt_path = None
    converted_pptx_path = None
    is_converted = False
    original_path = fp
    
    if fp.suffix.lower() == ".pptx":
        logger.info(f"[PPTX] PPTX 파일 감지, PPT로 변환 시작: {fp.name} (경로: {fp})")
        if not fp.exists():
            logger.error(f"[PPTX] 원본 PPTX 파일이 존재하지 않습니다: {fp}")
            return "", []
        
        converted_ppt_path = _convert_pptx_to_ppt(fp)
        if converted_ppt_path and converted_ppt_path.exists():
            file_size = converted_ppt_path.stat().st_size
            logger.info(f"[PPTX] PPTX->PPT 변환 성공: {converted_ppt_path.name} (크기: {file_size} bytes)")
            fp = converted_ppt_path
            is_converted = True
        else:
            logger.error(f"[PPTX] PPTX 파일 변환 실패: {original_path.name}. 변환된 파일이 없습니다.")
            logger.error(f"[PPTX] 원본 파일 경로: {original_path}, 존재 여부: {original_path.exists()}")
            return "", []
    
    try:
        from pptx import Presentation
    except ImportError:
        logger.warning(f"python-pptx not available, cannot process PPT/PPTX. Returning empty for {fp.name}")
        # 변환된 파일이 있으면 삭제
        if is_converted and converted_ppt_path and converted_ppt_path.exists():
            try:
                converted_ppt_path.unlink()
                logger.debug(f"[PPTX] 임시 변환 PPT 파일 삭제: {converted_ppt_path.name}")
            except Exception as e:
                logger.warning(f"[PPTX] 임시 변환 PPT 파일 삭제 실패: {converted_ppt_path.name}, 오류: {e}")
        return "", []
    
    try:
        # PPT 파일을 PPTX로 변환 (python-pptx는 PPTX만 지원)
        if fp.suffix.lower() == ".ppt":
            logger.info(f"[PPTX] PPT 파일을 PPTX로 변환하여 파싱: {fp.name}")
            converted_pptx_path = _convert_ppt_to_pptx_for_parsing(fp)
            if not converted_pptx_path or not converted_pptx_path.exists():
                logger.error(f"[PPTX] PPT를 PPTX로 변환 실패: {fp.name}")
                return "", []
            # PPTX 파일로 교체하여 파싱
            fp = converted_pptx_path
        
        # PPTX 파일 파싱
        prs = Presentation(str(fp))
        texts = []
        tables = []
        for i, slide in enumerate(prs.slides, start=1):
            for sh in slide.shapes:
                if hasattr(sh, "has_text_frame") and sh.has_text_frame:
                    txt = "\n".join(p.text for p in sh.text_frame.paragraphs if p.text)
                    txt = _clean_text(txt)
                    if txt:
                        texts.append(txt)
                if hasattr(sh, "has_table") and sh.has_table:
                    rows = []
                    for r in sh.table.rows:
                        rows.append([_clean_text(c.text) for c in r.cells])
                    if rows:
                        md = "\n".join("| " + " | ".join(row) + " |" for row in rows)
                        tables.append({"page": i, "bbox": [], "text": md})
        
        result_text = _clean_text("\n\n".join(texts))
        
        if not result_text.strip() and not tables:
            logger.warning(f"PPTX extraction may have failed for {original_path.name}, extracted text is too short or empty")
        
        return result_text, tables
        
    except Exception as e:
        logger.exception(f"Failed to extract PPTX {original_path.name}: {e}")
        return "", []
    finally:
        # 변환된 임시 PPT 파일 삭제
        if is_converted and converted_ppt_path and converted_ppt_path.exists():
            try:
                converted_ppt_path.unlink()
                logger.debug(f"[PPTX] 임시 변환 PPT 파일 삭제: {converted_ppt_path.name}")
            except Exception as e:
                logger.warning(f"[PPTX] 임시 변환 PPT 파일 삭제 실패: {converted_ppt_path.name}, 오류: {e}")
        
        # 임시 PPTX 파일 삭제 (PPT를 PPTX로 변환한 경우)
        if converted_pptx_path and converted_pptx_path.exists():
            try:
                converted_pptx_path.unlink()
                logger.debug(f"[PPTX] 임시 변환 PPTX 파일 삭제: {converted_pptx_path.name}")
            except Exception as e:
                logger.warning(f"[PPTX] 임시 변환 PPTX 파일 삭제 실패: {converted_pptx_path.name}, 오류: {e}")


def _convert_ppt_to_pptx_for_parsing(src: Path) -> Path | None:
    """PPT 파일을 PPTX로 변환 (python-pptx 파싱을 위해)"""
    try:
        import subprocess
        import shutil
        import time
        
        # LibreOffice 명령어 확인
        libreoffice_cmd = shutil.which("libreoffice")
        if not libreoffice_cmd:
            logger.warning(f"[PPT->PPTX] LibreOffice가 설치되지 않았습니다: {src.name}")
            return None
        
        outdir = src.parent
        logger.debug(f"[PPT->PPTX] 변환 시작 (파싱용): {src.name}")
        
        result = subprocess.run(
            [
                libreoffice_cmd, "--headless", "--convert-to", "pptx",
                "--outdir", str(outdir), str(src)
            ],
            check=False,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            timeout=120,
        )
        
        if result.returncode != 0:
            stderr_msg = result.stderr.decode("utf-8", errors="ignore") if result.stderr else ""
            logger.warning(f"[PPT->PPTX] 변환 실패: {src.name}. 에러: {stderr_msg[:200]}")
            return None
        
        # 변환된 파일 찾기
        cand = src.with_suffix(".pptx")
        max_wait = 10
        wait_interval = 0.5
        waited = 0
        while not cand.exists() and waited < max_wait:
            time.sleep(wait_interval)
            waited += wait_interval
        
        if cand.exists():
            file_size = cand.stat().st_size
            logger.info(f"[PPT->PPTX] 변환 성공 (파싱용): {cand.name} (크기: {file_size} bytes)")
            return cand
        
        # 대소문자 변형 확인
        parent = src.parent
        stem = src.stem
        for name in [f"{stem}.pptx", f"{stem}.PPTX", f"{stem}.Pptx"]:
            alt_path = parent / name
            if alt_path.exists():
                file_size = alt_path.stat().st_size
                logger.info(f"[PPT->PPTX] 변환 성공 (대체 경로, 파싱용): {alt_path.name} (크기: {file_size} bytes)")
                return alt_path
        
        logger.error(f"[PPT->PPTX] 변환된 파일을 찾을 수 없습니다: {src.name}")
        return None
        
    except subprocess.TimeoutExpired:
        logger.error(f"[PPT->PPTX] 변환 타임아웃: {src.name} (120초 초과)")
        return None
    except Exception as e:
        logger.exception(f"[PPT->PPTX] 변환 중 예외: {src.name}, 오류: {e}")
        return None

