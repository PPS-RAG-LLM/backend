"""
RAG 전처리 라우터 모듈
확장자별로 적절한 전처리 함수를 호출하는 라우팅 기능 제공
"""
from __future__ import annotations
import logging
from pathlib import Path
from typing import Optional

logger = logging.getLogger(__name__)


def _ext(p: Path) -> str:
    """파일 확장자 반환 (소문자)"""
    return p.suffix.lower()


def _clean_text(s: str | None) -> str:
    """텍스트 정규화 (pdf_preprocessing에서 import)"""
    from service.preprocessing.extension.pdf_preprocessing import _clean_text as _clean_text_pdf
    return _clean_text_pdf(s)


def _extract_plain_text(fp: Path) -> tuple[str, list[dict]]:
    """TXT/MD 파일 추출"""
    try:
        text = fp.read_text(encoding="utf-8", errors="ignore")
    except Exception:
        text = fp.read_text(errors="ignore")
    return _clean_text(text), []


def _extract_docx(fp: Path) -> tuple[str, list[dict]]:
    """DOCX 파일 추출"""
    try:
        from docx import Document
    except Exception:
        logger.warning(f"python-docx not available, treating {fp} as plain text")
        return _extract_plain_text(fp)
    
    try:
        d = Document(str(fp))
        paras = [_clean_text(p.text) for p in d.paragraphs if _clean_text(p.text)]
        tables = []
        for tb in d.tables:
            rows = []
            for r in tb.rows:
                rows.append([_clean_text(c.text) for c in r.cells])
            if rows:
                md = "\n".join("| " + " | ".join(row) + " |" for row in rows)
                tables.append({"page": 0, "bbox": [], "text": md})
        return _clean_text("\n\n".join(paras)), tables
    except Exception as e:
        logger.exception(f"Failed to extract DOCX {fp}: {e}")
        return _extract_plain_text(fp)


def _extract_pptx(fp: Path) -> tuple[str, list[dict]]:
    """PPTX 파일 추출"""
    try:
        from pptx import Presentation
    except Exception:
        logger.warning(f"python-pptx not available, skipping {fp}")
        return "", []
    
    try:
        prs = Presentation(str(fp))
        texts = []
        tables = []
        for i, slide in enumerate(prs.slides, start=1):
            for sh in slide.shapes:
                if hasattr(sh, "has_text_frame") and sh.has_text_frame:
                    txt = "\n".join(p.text for p in sh.text_frame.paragraphs if p.text)
                    txt = _clean_text(txt)
                    if txt:
                        texts.append(txt)
                if hasattr(sh, "has_table") and sh.has_table:
                    rows = []
                    for r in sh.table.rows:
                        rows.append([_clean_text(c.text) for c in r.cells])
                    if rows:
                        md = "\n".join("| " + " | ".join(row) + " |" for row in rows)
                        tables.append({"page": i, "bbox": [], "text": md})
        return _clean_text("\n\n".join(texts)), tables
    except Exception as e:
        logger.exception(f"Failed to extract PPTX {fp}: {e}")
        return "", []


def _df_to_markdown(df, max_rows=500) -> str:
    """DataFrame을 마크다운 테이블로 변환"""
    if len(df) > max_rows:
        df = df.head(max_rows)
    cols = [str(c) for c in df.columns]
    lines = ["| " + " | ".join(cols) + " |"]
    for _, row in df.iterrows():
        lines.append("| " + " | ".join(_clean_text(str(v)) for v in row.tolist()) + " |")
    return "\n".join(lines)


def _extract_csv(fp: Path) -> tuple[str, list[dict]]:
    """CSV 파일 추출"""
    try:
        import pandas as pd
        df = pd.read_csv(fp)
        md = _df_to_markdown(df)
        return "", [{"page": 0, "bbox": [], "text": md}]
    except Exception as e:
        logger.warning(f"Failed to extract CSV {fp}: {e}, trying as plain text")
        return _extract_plain_text(fp)


def _extract_excel(fp: Path) -> tuple[str, list[dict]]:
    """Excel 파일 추출"""
    try:
        import pandas as pd
        xls = pd.ExcelFile(fp)
        tables = []
        for name in xls.sheet_names:
            df = xls.parse(name)
            md = f"### {name}\n" + _df_to_markdown(df)
            tables.append({"page": 0, "bbox": [], "text": md})
        return "", tables
    except Exception as e:
        logger.warning(f"Failed to extract Excel {fp}: {e}")
        return "", []


def _convert_via_libreoffice(src: Path, target_ext: str) -> Optional[Path]:
    """LibreOffice를 통한 문서 변환"""
    try:
        import subprocess
        outdir = src.parent
        subprocess.run([
            "libreoffice", "--headless", "--convert-to", target_ext, 
            "--outdir", str(outdir), str(src)
        ], check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
        cand = src.with_suffix("." + target_ext)
        return cand if cand.exists() else None
    except Exception:
        return None


def _extract_doc(fp: Path) -> tuple[str, list[dict]]:
    """DOC 파일 추출 (LibreOffice 변환 시도)"""
    conv = _convert_via_libreoffice(fp, "docx")
    if conv and conv.exists():
        result = _extract_docx(conv)
        try:
            conv.unlink()  # 임시 변환 파일 삭제
        except Exception:
            pass
        return result
    return _extract_plain_text(fp)


def _extract_ppt(fp: Path) -> tuple[str, list[dict]]:
    """PPT 파일 추출 (LibreOffice 변환 시도)"""
    conv = _convert_via_libreoffice(fp, "pptx")
    if conv and conv.exists():
        result = _extract_pptx(conv)
        try:
            conv.unlink()  # 임시 변환 파일 삭제
        except Exception:
            pass
        return result
    return "", []


def _extract_any(path: Path) -> tuple[str, list[dict]]:
    """통합 문서 추출 라우터"""
    from service.preprocessing.extension.pdf_preprocessing import _extract_pdf_with_tables
    from service.preprocessing.extension.hwp_preprocessing import _extract_hwp
    
    ext = _ext(path)
    if ext == ".pdf":
        # PDF는 페이지 정보를 포함하지만, _extract_any는 (text, tables)만 반환
        text, tables, _, _ = _extract_pdf_with_tables(path)
        return text, tables
    if ext in {".txt", ".text", ".md"}:
        return _extract_plain_text(path)
    if ext == ".docx":
        return _extract_docx(path)
    if ext == ".pptx":
        return _extract_pptx(path)
    if ext == ".csv":
        return _extract_csv(path)
    if ext in {".xlsx", ".xls"}:
        return _extract_excel(path)
    if ext == ".doc":
        return _extract_doc(path)
    if ext == ".ppt":
        return _extract_ppt(path)
    if ext == ".hwp":
        return _extract_hwp(path)
    # 모르는 확장자는 텍스트로 시도
    return _extract_plain_text(path)


async def extract_documents():
    """
    문서 전처리 메인 함수
    모든 확장자의 파일을 _extract_any를 사용하여 처리합니다.
    
    Returns:
        dict: 전처리 결과
    """
    import json
    import shutil
    from collections import defaultdict
    from tqdm import tqdm  # type: ignore
    
    # 필요한 상수 및 함수 import
    from service.admin.manage_vator_DB import (
        EXTRACTED_TEXT_DIR,
        LOCAL_DATA_ROOT,
        RAW_DATA_DIR,
        META_JSON_PATH,
        TASK_TYPES,
        SUPPORTED_EXTS,
        _determine_level_for_task,
        _parse_doc_version,
        get_security_level_rules_all,
    )
    from service.preprocessing.extension.pdf_preprocessing import _extract_pdf_with_tables, _clean_text
    from utils.time import now_kst_string
    
    EXTRACTED_TEXT_DIR.mkdir(parents=True, exist_ok=True)
    LOCAL_DATA_ROOT.mkdir(parents=True, exist_ok=True)
    RAW_DATA_DIR.mkdir(parents=True, exist_ok=True)

    # 규칙 로드
    all_rules = get_security_level_rules_all()  # {task: {"maxLevel":N, "levels":{...}}}

    # 이전 메타 로드
    prev_meta: dict[str, dict] = {}
    if META_JSON_PATH.exists():
        try:
            prev_meta = json.loads(META_JSON_PATH.read_text(encoding="utf-8"))
        except Exception:
            logger.exception("Failed to read META JSON; recreating.")

    # 중복 제거 로직: 파일명 마지막 토큰이 날짜/버전 숫자면 최신만 유지
    def _extract_base_and_date(p: Path):
        name = p.stem
        parts = name.split("_")
        date_num = 0
        if len(parts) >= 2:
            cand = parts[-1]
            if cand.isdigit() and len(cand) in (4, 6, 8):
                try:
                    date_num = int(cand)
                except Exception:
                    date_num = 0
        mid_tokens = [t for t in parts[:-1] if t and not t.isdigit()]
        base = max(mid_tokens, key=len) if mid_tokens else parts[0]
        return base, date_num

    raw_files = [p for p in RAW_DATA_DIR.rglob("*") if p.is_file() and _ext(p) in SUPPORTED_EXTS]

    # base(문서ID 유사)별로 버전 후보 묶기: (Path, date_num)
    grouped: dict[str, list[tuple[Path, int]]] = defaultdict(list)
    for p in raw_files:
        base, date_num = _extract_base_and_date(p)
        grouped[base].append((p, date_num))

    kept, removed = [], []
    for base, lst in grouped.items():
        # lst 원소는 (Path, date_num)
        lst_sorted = sorted(
            lst,
            key=lambda it: (it[1], it[0].stat().st_mtime, len(it[0].name))  # date_num → mtime → 이름 길이
        )
        keep_path = lst_sorted[-1][0]          # 최신 후보의 Path
        kept.append(keep_path)
        for old_path, _ in lst_sorted[:-1]:
            try:
                old_path.unlink(missing_ok=True)
                removed.append(str(old_path.relative_to(RAW_DATA_DIR)))
            except Exception:
                logger.exception("Failed to remove duplicate: %s", old_path)

    if not kept:
        return {
            "message": "처리할 문서가 없습니다.",
            "meta_path": str(META_JSON_PATH),
            "deduplicated": {"removedCount": len(removed), "removed": removed},
        }

    new_meta: dict[str, dict] = {}
    for src in tqdm(kept, desc="문서 전처리"):
        try:
            logger.info(f"[Extract] 파일 처리 시작: {src.name}")
            
            # _extract_any를 사용하여 모든 확장자 처리
            pages_text_dict: dict[int, str] = {}
            total_pages = 0
            
            # PDF인 경우 페이지별 정보를 포함하여 추출
            if _ext(src) == ".pdf":
                text, tables, pages_text_dict, total_pages = _extract_pdf_with_tables(src)
            else:
                # PDF가 아닌 파일은 _extract_any 사용
                text, tables = _extract_any(src)
            
            # 표 추출 결과 로깅
            if tables:
                logger.info(f"[Extract] {src.name}: 표 {len(tables)}개 추출됨")
                for t_idx, t in enumerate(tables):
                    page = t.get("page", 0)
                    text_preview = (t.get("text", "")[:100] + "...") if t.get("text") else ""
                    logger.info(f"[Extract] {src.name} 표 {t_idx+1}: 페이지={page}, 텍스트 미리보기={text_preview}")
            else:
                logger.info(f"[Extract] {src.name}: 추출된 표 없음 (텍스트만 추출됨)")
            
            # 작업유형별 보안 레벨 (본문+표 모두 포함해서 판정)
            whole_for_level = text + "\n\n" + "\n\n".join(t.get("text","") for t in (tables or []))
            sec_map = {task: _determine_level_for_task(
                whole_for_level, all_rules.get(task, {"maxLevel": 1, "levels": {}})
            ) for task in TASK_TYPES}

            max_sec = max(sec_map.values()) if sec_map else 1
            sec_folder = f"securityLevel{int(max_sec)}"

            rel_from_raw = src.relative_to(RAW_DATA_DIR)
            # 원본 그대로 보관
            dest_rel = Path(sec_folder) / rel_from_raw
            dest_abs = LOCAL_DATA_ROOT / dest_rel
            dest_abs.parent.mkdir(parents=True, exist_ok=True)
            try:
                shutil.copy2(src, dest_abs)
            except Exception:
                logger.exception("Failed to copy file: %s", dest_abs)

            # 통합 파일 저장 (.txt, 마크다운 형식) - 페이지 순서대로 텍스트와 표 배치
            # 파일명은 원본 파일명 그대로 사용 (확장자만 .txt)
            txt_rel = dest_rel.with_suffix(".txt")
            (EXTRACTED_TEXT_DIR / txt_rel).parent.mkdir(parents=True, exist_ok=True)
            saved_files: dict[str, str] = {}
            
            # 통합 파일 저장
            combined_txt_file = EXTRACTED_TEXT_DIR / txt_rel
            try:
                # 페이지별 표 그룹화
                pages_tables: dict[int, list[dict]] = defaultdict(list)
                for t in (tables or []):
                    page_num = t.get("page", 0)
                    if page_num > 0:
                        pages_tables[page_num].append(t)
                
                # 통합 파일 작성 (페이지 순서대로) - 페이지 마커 없이 순수 텍스트+표만 저장
                with open(combined_txt_file, "w", encoding="utf-8") as f:
                    # PDF인 경우: 페이지별 텍스트 정보 사용
                    if pages_text_dict:
                        # 모든 페이지 번호 수집 (텍스트와 표 모두 고려)
                        all_page_nums = set(pages_text_dict.keys())
                        all_page_nums.update(pages_tables.keys())
                        
                        if all_page_nums:
                            for page_num in sorted(all_page_nums):
                                # 해당 페이지의 텍스트 (페이지 마커 없이)
                                page_text_content = pages_text_dict.get(page_num, "")
                                if page_text_content:
                                    f.write(page_text_content)
                                    f.write("\n\n")
                                
                                # 해당 페이지의 표들 (텍스트 뒤에 삽입)
                                page_tables_list = pages_tables.get(page_num, [])
                                if page_tables_list:
                                    for t_idx, t in enumerate(page_tables_list):
                                        table_text = t.get("text", "")
                                        if table_text:
                                            f.write(table_text)
                                            f.write("\n\n")
                                
                                # 페이지 구분선 (마지막 페이지가 아니면)
                                if page_num < max(all_page_nums):
                                    f.write("\n---\n\n")
                        else:
                            # 페이지 정보가 없으면 전체 텍스트만
                            if text.strip():
                                f.write(text)
                                f.write("\n\n")
                            if tables:
                                for t in tables:
                                    table_text = t.get("text", "")
                                    if table_text:
                                        f.write(table_text)
                                        f.write("\n\n")
                    else:
                        # PDF가 아닌 경우: 전체 텍스트와 표만 저장
                        if text.strip():
                            f.write(text)
                            f.write("\n\n")
                        if tables:
                            for t in tables:
                                table_text = t.get("text", "")
                                if table_text:
                                    f.write(table_text)
                                    f.write("\n\n")
                
                saved_files["text"] = str(combined_txt_file)
                logger.info(f"[Extract] 통합 파일 저장: {combined_txt_file}")
            except Exception as e:
                logger.exception(f"[Extract] 통합 파일 저장 실패: {e}")

            # doc_id/version 유추
            stem = rel_from_raw.stem
            doc_id, version = _parse_doc_version(stem)

            info = {
                "chars": len(text),
                "lines": len(text.splitlines()),
                "preview": (_clean_text(text[:200].replace("\n"," ")) + "…") if text else "",
                "security_levels": sec_map,  # 작업유형별 보안레벨
                "doc_id": doc_id,
                "version": version,
                "tables": tables or [],  # ★ 표 정보 추가
                "pages": pages_text_dict if pages_text_dict else {},  # ★ 페이지별 텍스트 정보 (메타데이터용)
                "total_pages": total_pages,  # ★ 총 페이지 수
                "sourceExt": _ext(src),  # 원본 확장자 기록
                "saved_files": saved_files,  # 저장된 파일 경로 추가
                # 페이로드 정보 (LLM에 전달하지 않지만 메타데이터로 저장)
                "extraction_info": {
                    "original_file": src.name,
                    "text_length": len(text),
                    "table_count": len(tables or []),
                    "extracted_at": now_kst_string(),
                },
            }
            new_meta[str(dest_rel)] = info
            logger.info(f"[Extract] {src.name}: 메타데이터 저장 완료 (텍스트={len(text)}자, 표={len(tables or [])}개)")

        except Exception as e:
            logger.exception("Failed to process: %s", src)
            try:
                rel_from_raw = src.relative_to(RAW_DATA_DIR)
                dest_rel = Path("securityLevel1") / rel_from_raw
                new_meta[str(dest_rel)] = {"error": str(e)}
            except Exception:
                new_meta[src.name] = {"error": str(e)}

    META_JSON_PATH.write_text(
        json.dumps(new_meta, ensure_ascii=False, indent=2), encoding="utf-8"
    )
    return {
        "message": "문서 추출 완료",
        "file_count": len(kept),
        "meta_path": str(META_JSON_PATH),
        "deduplicated": {"removedCount": len(removed), "removed": removed},
    }

