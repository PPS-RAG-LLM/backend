# === Vector DB Service (Milvus Server, Pro) ===
# - 작업유형(task_type)별 보안레벨 관리: doc_gen | summary | qna
# - Milvus Docker 서버 전용 (Lite 제거)
# - 벡터/하이브리드 검색 지원, 실행 로그 적재

from __future__ import annotations
import uuid
import json
import os
import time
import logging
import re
import unicodedata

# sqlite3 제거
import shutil
import threading
import asyncio
from pathlib import Path
from typing import List, Dict, Optional, Tuple, Any
from collections import defaultdict, Counter

from pydantic import BaseModel, Field
from pymilvus import MilvusClient, DataType

from config import config as app_config
from repository.embedding_model import get_active_embedding_model_name
from repository.rag_settings import get_vector_settings_row

try:
    # Milvus 2.4+ Function/BM25 하이브리드
    from pymilvus import Function, FunctionType
except Exception:
    Function = None
    class FunctionType:
        BM25 = "BM25"

# ORM 추가 임포트
from utils.database import get_session
from storage.db_models import (
    EmbeddingModel,
    RagSettings,
    SecurityLevelConfigTask,
    SecurityLevelKeywordsTask,
)

def _split_for_varchar_bytes(
    text: str,
    hard_max_bytes: int = 32768,
    soft_max_bytes: int = 30000,   # 여유 버퍼
    table_mark: str = "[[TABLE",
) -> list[str]:
    """
    VARCHAR 초과 방지: UTF-8 바이트 기준으로 안전 분할.
    - 표 텍스트는 헤더([[TABLE ...]])를 첫 조각에만 포함.
    - 이후 조각엔 [[TABLE_CONT i/n]] 마커를 부여.
    - 개행 경계 우선(backtrack), 그래도 안되면 하드컷.
    """
    if not text:
        return [""]

    # 표 헤더 분리
    header = ""
    body = text
    if text.startswith(table_mark):
        head_end = text.find("]]")
        if head_end != -1:
            head_end += 2
            if head_end < len(text) and text[head_end] == "\n":
                head_end += 1
            header, body = text[:head_end], text[head_end:]

    def _split_body(b: str) -> list[str]:
        out: list[str] = []
        b_bytes = b.encode("utf-8")
        n = len(b_bytes)
        i = 0
        while i < n:
            j = min(i + soft_max_bytes, n)
            # 개행 경계로 뒤로 물러나기
            k = j
            backtracked = False
            # j부터 i까지 역방향으로 \n 바이트(0x0A) 탐색
            while k > i and (j - k) < 2000:  # 최대 2KB만 백트랙
                if b_bytes[k-1:k] == b"\n":
                    backtracked = True
                    break
                k -= 1
            if backtracked and (k - i) >= int(soft_max_bytes * 0.6):
                cut = k
            else:
                cut = j

            # 하드 컷(멀티바이트 경계 맞추기)
            if cut - i > hard_max_bytes:
                cut = i + hard_max_bytes

            # UTF-8 안전 디코드: 경계가 문자를 반쯤 자를 수 있으니 넉넉히 조정
            chunk = b_bytes[i:cut]
            # 만약 디코드 에러가 나면 한 바이트씩 줄이며 안전 경계 찾기
            while True:
                try:
                    s = chunk.decode("utf-8")
                    break
                except UnicodeDecodeError:
                    cut -= 1
                    if cut <= i:
                        # 최악의 경우 한 글자라도 디코드되게 한 바이트 앞당김
                        cut = i + 1
                    chunk = b_bytes[i:cut]
            out.append(s)
            i = cut
        return out

    if len(text.encode("utf-8")) <= hard_max_bytes:
        return [text]

    parts = _split_body(body)
    if header:
        total = len(parts)
        result = []
        for idx, c in enumerate(parts, start=1):
            if idx == 1:
                # 첫 조각은 헤더 + 본문
                # 전체가 하드맥스를 넘지 않게 헤더와 합친 뒤 한번 더 자르기
                first = header + c
                if len(first.encode("utf-8")) <= hard_max_bytes:
                    result.append(first)
                else:
                    # 너무 크면 헤더는 유지하고 c를 다시 잘라 붙임
                    # (헤더가 길 때 매우 예외적)
                    subparts = _split_body(c)
                    if subparts:
                        # 첫 조각은 헤더 + 첫 sub
                        f = header + subparts[0]
                        if len(f.encode("utf-8")) > hard_max_bytes:
                            # 헤더 자체가 큰 극단: 헤더만 넣고 이후 CONT로 처리
                            result.append(header[:0] + header)  # 그대로
                            # 나머지는 CONT
                            for sidx, sp in enumerate(subparts, start=1):
                                tag = f"[[TABLE_CONT {sidx}/{len(subparts)}]]\n"
                                result.append(tag + sp)
                        else:
                            result.append(f)
                            # 나머지는 CONT
                            for sidx, sp in enumerate(subparts[1:], start=2):
                                tag = f"[[TABLE_CONT {sidx}/{len(subparts)}]]\n"
                                result.append(tag + sp)
                    else:
                        result.append(header)  # 본문이 없으면 헤더만
            else:
                tag = f"[[TABLE_CONT {idx}/{total}]]\n"
                # tag + c 가 하드맥스를 넘지 않도록 재자르기
                rest = tag + c
                if len(rest.encode("utf-8")) <= hard_max_bytes:
                    result.append(rest)
                else:
                    subs = _split_body(c)
                    for sidx, sp in enumerate(subs, start=1):
                        subt = f"[[TABLE_CONT {idx}.{sidx}/{total}]]\n" + sp
                        if len(subt.encode("utf-8")) <= hard_max_bytes:
                            result.append(subt)
                        else:
                            # 그래도 넘으면 하드컷으로 마지막 방어
                            bb = subt.encode("utf-8")[:hard_max_bytes]
                            result.append(bb.decode("utf-8", errors="ignore"))
        return result
    else:
        return parts


# KST 시간 포맷 유틸
from utils.time import now_kst, now_kst_string

from service.retrieval.common import get_or_load_hf_embedder, hf_embed_text, chunk_text
from service.retrieval.reranker import rerank_snippets
from utils.model_load import resolve_model_input

logger = logging.getLogger(__name__)

# -------------------------------------------------
# 경로 상수
# -------------------------------------------------
BASE_DIR = Path(__file__).resolve().parent  # .../backend/service/admin
PROJECT_ROOT = BASE_DIR.parent.parent  # .../backend
_RETRIEVAL_CFG: Dict[str, Any] = app_config.get("retrieval", {}) or {}
_RETRIEVAL_PATHS: Dict[str, str] = _RETRIEVAL_CFG.get("paths", {}) or {}
_VECTOR_DEFAULTS: Dict[str, Any] = app_config.get("vector_defaults", {}) or {}


def _cfg_path(key: str, fallback: str) -> Path:
    value = _RETRIEVAL_PATHS.get(key, fallback)
    return (PROJECT_ROOT / Path(value)).resolve()


STORAGE_DIR = _cfg_path("storage_dir", "storage")
USER_DATA_ROOT = _cfg_path("user_data_root", "storage/user_data")
RAW_DATA_DIR = _cfg_path("raw_data_dir", "storage/user_data/row_data")
LOCAL_DATA_ROOT = _cfg_path("local_data_root", "storage/user_data/preprocessed_data")
RESOURCE_DIR = _cfg_path("resources_dir", str(BASE_DIR / "resources"))
EXTRACTED_TEXT_DIR = _cfg_path("extracted_text_dir", "storage/extracted_texts")
META_JSON_PATH = _cfg_path("meta_json_path", "storage/extracted_texts/_extraction_meta.json")
MODEL_ROOT_DIR = _cfg_path("model_root_dir", "storage/embedding-models")
RERANK_MODEL_PATH = _cfg_path("rerank_model_path", "storage/rerank_model/Qwen3-Reranker-0.6B")
VAL_SESSION_ROOT = _cfg_path("val_session_root", "storage/val_data")
SESSIONS_INDEX_PATH = _cfg_path("sessions_index_path", "storage/val_data/_sessions.json")
DATABASE_CFG = app_config.get("database", {}) or {}
SQLITE_DB_PATH = (PROJECT_ROOT / Path(DATABASE_CFG.get("path", "storage/pps_rag.db"))).resolve()

_MILVUS_CFG = _RETRIEVAL_CFG.get("milvus", {}) or {}
MILVUS_URI = os.getenv("MILVUS_URI", _MILVUS_CFG.get("uri", "http://localhost:19530"))
MILVUS_TOKEN = os.getenv("MILVUS_TOKEN", _MILVUS_CFG.get("token") or None)
COLLECTION_NAME = _MILVUS_CFG.get("collection", "pdf_chunks_pro")
TASK_TYPES = tuple(_RETRIEVAL_CFG.get("task_types") or ("doc_gen", "summary", "qna"))
SUPPORTED_EXTS = set(_RETRIEVAL_CFG.get("supported_extensions"))

ZERO_WIDTH_RE = re.compile(r"[\u200B-\u200D\u2060\uFEFF]")
CONTROL_RE = re.compile(r"[\x00-\x08\x0B\x0C\x0E-\x1F]")
MULTISPACE_LINE_END_RE = re.compile(r"[ \t]+\n")
NEWLINES_RE = re.compile(r"\n{3,}")

_CURRENT_EMBED_MODEL_KEY = str(_RETRIEVAL_CFG.get("default_embedding_key", "qwen3_0_6b"))
_CURRENT_SEARCH_TYPE = str(_RETRIEVAL_CFG.get("search_type", "hybrid"))
_CURRENT_CHUNK_SIZE = int(_VECTOR_DEFAULTS.get("chunk_size", 512))
_CURRENT_OVERLAP = int(_VECTOR_DEFAULTS.get("overlap", 64))


# -------------------------------------------------
# 텍스트 정리 및 다중 확장자 지원
# -------------------------------------------------

def _clean_text(s: str | None) -> str:
    """PDF 특수문자 제거 및 텍스트 정규화"""
    if not s:
        return ""
    s = unicodedata.normalize("NFKC", s)
    # 흔한 공백/불릿/문장부호 정리
    s = s.replace("\xa0", " ").replace("\u2022", "•").replace("\u2212", "-")
    s = ZERO_WIDTH_RE.sub("", s)
    s = CONTROL_RE.sub("", s)
    s = MULTISPACE_LINE_END_RE.sub("\n", s)
    s = NEWLINES_RE.sub("\n\n", s)
    return s.strip()


def _ext(p: Path) -> str:
    """파일 확장자 반환 (소문자)"""
    return p.suffix.lower()

def _df_to_markdown_repeat_header(df, max_rows=500) -> str:
    """
    각 데이터 행 앞에 헤더를 반복해서 출력하는 마크다운 테이블 생성.
    예)
      이름 | 나이 | 성
      기정 | 29  | 남자
      이름 | 나이 | 성
      수현 | 26  | 여자
    """
    import pandas as pd
    if isinstance(df, pd.DataFrame) and len(df) > max_rows:
        df = df.head(max_rows)

    cols = [str(c) for c in df.columns]
    header = "| " + " | ".join(cols) + " |"

    lines = []
    for _, row in df.iterrows():
        lines.append(header)
        lines.append("| " + " | ".join(_clean_text(str(v)) for v in row.tolist()) + " |")
    return "\n".join(lines)


def _markdown_repeat_header(md: str) -> str:
    """기존 마크다운 표 문자열에서 데이터 행 앞에 헤더를 반복해 반환."""
    if not md:
        return md

    lines = [line for line in md.splitlines() if line.strip()]
    if len(lines) <= 1:
        return md

    header = lines[0]
    sep_re = re.compile(r"^\s*\|?(?:\s*:?-+:?\s*\|)+\s*$")
    data_lines: list[str] = []
    for line in lines[1:]:
        if sep_re.match(line):
            continue
        data_lines.append(line)

    if not data_lines:
        return md

    out: list[str] = []
    for dl in data_lines:
        out.append(header)
        out.append(dl)
    return "\n".join(out)


def _extract_pdf_with_tables(pdf_path: Path) -> tuple[str, list[dict], Dict[int, str], int]:
    """
    PDF에서 본문 텍스트는 PyMuPDF로 추출하고,
    표는 PyMuPDF find_tables()와 Tabula로 추출해 마크다운 테이블로 반환한다.
    - 페이지별로 텍스트와 표를 분리 추출
    - 페이지 정보는 메타데이터에만 저장 (텍스트에는 페이지 마커 없음)
    - 이미지는 처리하지 않음
    - Tabula가 없거나 실패 시: PyMuPDF find_tables()로 폴백
    
    Returns:
        tuple: (전체 텍스트, 표 리스트, 페이지별 텍스트 딕셔너리, 총 페이지 수)
    """
    import fitz  # PyMuPDF
    
    page_texts: list[str] = []
    all_tables: list[dict] = []
    total_pages = 0
    
    # with 컨텍스트 매니저로 안전하게 문서 열기/닫기
    try:
        with fitz.open(pdf_path) as doc:
            # 페이지 수를 먼저 저장 (doc이 닫히기 전에)
            total_pages = doc.page_count
            
            for page_idx in range(total_pages):
                page_num = page_idx + 1
                page = doc[page_idx]
                page_text = ""
                page_tables: list[dict] = []
                
                # 1. 일반 텍스트 추출 (PyMuPDF)
                try:
                    raw_text = page.get_text()
                    if raw_text and raw_text.strip():
                        clean_text = _clean_text(raw_text)
                        # 페이지 마커 없이 순수 텍스트만 저장 (페이지 정보는 메타데이터에 저장)
                        page_texts.append((page_num, clean_text))  # (페이지번호, 텍스트) 튜플로 저장
                except Exception as e:
                    logger.warning(f"[PyMuPDF] 텍스트 추출 실패 (p{page_num}): {e}")
                
                # 2. 표 추출 (PyMuPDF 내장 테이블 감지)
                try:
                    tables = page.find_tables()
                    table_list = getattr(tables, "tables", tables) if tables else []
                    
                    if table_list:
                        logger.info(f"[PyMuPDF] {page_num}페이지: 테이블 {len(table_list)}개 탐지")
                    
                    for i, table in enumerate(table_list):
                        try:
                            table_data = table.extract()
                            if not table_data:
                                logger.debug(f"[PyMuPDF] {page_num}페이지 테이블 {i+1}: 빈 데이터")
                                continue
                            
                            # 마크다운 테이블 생성
                            markdown_lines: list[str] = []
                            
                            # 헤더
                            header_row = [_clean_text(str(cell or "")) for cell in table_data[0]]
                            markdown_lines.append("| " + " | ".join(header_row) + " |")
                            markdown_lines.append("|" + "---|" * len(header_row))
                            
                            # 데이터 행
                            for row in table_data[1:]:
                                clean_row = [_clean_text(str(cell or "")) for cell in row]
                                markdown_lines.append("| " + " | ".join(clean_row) + " |")
                            
                            md = "\n".join(markdown_lines)
                            md = _markdown_repeat_header(md)  # 헤더 반복 적용
                            
                            if md.strip():
                                rect = fitz.Rect(*table.bbox) if hasattr(table, 'bbox') else None
                                page_tables.append({
                                    "page": page_num,
                                    "bbox": [float(rect.x0), float(rect.y0), float(rect.x1), float(rect.y1)] if rect else [],
                                    "text": md,
                                })
                                logger.info(f"[PyMuPDF] {page_num}페이지 테이블 {i+1}: 마크다운 변환 성공, 길이={len(md)}")
                        except Exception as e:
                            logger.exception(f"[PyMuPDF] {page_num}페이지 테이블 {i+1} 처리 실패: {e}")
                except Exception as e:
                    logger.warning(f"[PyMuPDF] {page_num}페이지 표 감지 실패: {e}")
                
                # 3. Tabula로 추가 표 추출 (선택)
                try:
                    import tabula  # requires Java (JRE/JDK)
                    TABULA_AVAILABLE = True
                except ImportError:
                    TABULA_AVAILABLE = False
                
                if TABULA_AVAILABLE:
                    try:
                        dfs = tabula.read_pdf(
                            str(pdf_path),
                            pages=page_num,
                            multiple_tables=True,
                            lattice=True,
                            silent=True,
                        )
                        if not dfs:
                            dfs = tabula.read_pdf(
                                str(pdf_path),
                                pages=page_num,
                                multiple_tables=True,
                                stream=True,
                                silent=True,
                            )
                        
                        if not dfs:
                            logger.debug(f"[Tabula] {page_num}페이지: 탐지된 테이블 없음")
                        else:
                            logger.info(f"[Tabula] {page_num}페이지: 테이블 {len(dfs)}개 탐지, shapes={[df.shape for df in dfs]}")
                        
                        for j, df in enumerate(dfs or []):
                            try:
                                if df.empty:
                                    logger.debug(f"[Tabula] {page_num}페이지 테이블 {j+1}: empty DataFrame - 건너뜀")
                                    continue
                                
                                logger.info(f"[Tabula] {page_num}페이지 테이블 {j+1}: columns={list(df.columns)}, shape={df.shape}")
                                
                                md = _df_to_markdown_repeat_header(df)
                                md = _clean_text(md)
                                
                                if md:
                                    # PyMuPDF에서 이미 같은 페이지의 표를 찾았는지 확인 (중복 제거)
                                    is_duplicate = False
                                    for existing_table in page_tables:
                                        if existing_table.get("page") == page_num:
                                            # 간단한 중복 체크: 첫 줄이 비슷하면 중복으로 간주
                                            existing_first_line = existing_table.get("text", "").split("\n")[0] if existing_table.get("text") else ""
                                            new_first_line = md.split("\n")[0] if md else ""
                                            if existing_first_line and new_first_line and existing_first_line[:50] == new_first_line[:50]:
                                                is_duplicate = True
                                                logger.debug(f"[Tabula] {page_num}페이지 테이블 {j+1}: PyMuPDF와 중복으로 판단, 건너뜀")
                                                break
                                    
                                    if not is_duplicate:
                                        page_tables.append({
                                            "page": page_num,
                                            "bbox": [],
                                            "text": md,
                                        })
                                        logger.info(f"[Tabula] {page_num}페이지 테이블 {j+1}: 마크다운 변환 성공, 길이={len(md)}")
                                else:
                                    logger.debug(f"[Tabula] {page_num}페이지 테이블 {j+1}: 마크다운 변환 후 빈 문자열")
                            except Exception as e:
                                logger.exception(f"[Tabula] {page_num}페이지 테이블 {j+1} 처리 실패: {e}")
                                continue
                    except Exception as e:
                        logger.debug(f"[Tabula] {page_num}페이지 표 추출 실패: {e}")
                
                # 페이지별 표를 전체 리스트에 추가
                all_tables.extend(page_tables)
    
    except Exception:
        logger.exception(f"Failed to open PDF: {pdf_path}")
        return "", []  # 열리지 않는 PDF는 빈 결과 반환
    
    # 모든 페이지 텍스트 결합 (페이지 마커 없이)
    pdf_text = _clean_text("\n\n".join(text for _, text in page_texts if text))
    
    # 페이지별 텍스트 딕셔너리 생성 (메타데이터용)
    pages_text_dict: Dict[int, str] = {}
    for page_num, text_content in page_texts:
        if text_content:
            pages_text_dict[page_num] = text_content

    # 표 추출 결과 로깅
    if all_tables:
        logger.info(f"[Extract] {pdf_path.name}: 총 {len(all_tables)}개 표 추출 완료")
    else:
        logger.info(f"[Extract] {pdf_path.name}: 추출된 표 없음")
    
    return pdf_text, all_tables, pages_text_dict, total_pages


def _extract_plain_text(fp: Path) -> tuple[str, list[dict]]:
    """TXT/MD 파일 추출"""
    try:
        text = fp.read_text(encoding="utf-8", errors="ignore")
    except Exception:
        text = fp.read_text(errors="ignore")
    return _clean_text(text), []


def _extract_docx(fp: Path) -> tuple[str, list[dict]]:
    """DOCX 파일 추출"""
    try:
        from docx import Document
    except Exception:
        logger.warning(f"python-docx not available, treating {fp} as plain text")
        return _extract_plain_text(fp)
    
    try:
        d = Document(str(fp))
        paras = [_clean_text(p.text) for p in d.paragraphs if _clean_text(p.text)]
        tables = []
        for tb in d.tables:
            rows = []
            for r in tb.rows:
                rows.append([_clean_text(c.text) for c in r.cells])
            if rows:
                md = "\n".join("| " + " | ".join(row) + " |" for row in rows)
                tables.append({"page": 0, "bbox": [], "text": md})
        return _clean_text("\n\n".join(paras)), tables
    except Exception as e:
        logger.exception(f"Failed to extract DOCX {fp}: {e}")
        return _extract_plain_text(fp)


def _extract_pptx(fp: Path) -> tuple[str, list[dict]]:
    """PPTX 파일 추출"""
    try:
        from pptx import Presentation
    except Exception:
        logger.warning(f"python-pptx not available, skipping {fp}")
        return "", []
    
    try:
        prs = Presentation(str(fp))
        texts = []
        tables = []
        for i, slide in enumerate(prs.slides, start=1):
            for sh in slide.shapes:
                if hasattr(sh, "has_text_frame") and sh.has_text_frame:
                    txt = "\n".join(p.text for p in sh.text_frame.paragraphs if p.text)
                    txt = _clean_text(txt)
                    if txt:
                        texts.append(txt)
                if hasattr(sh, "has_table") and sh.has_table:
                    rows = []
                    for r in sh.table.rows:
                        rows.append([_clean_text(c.text) for c in r.cells])
                    if rows:
                        md = "\n".join("| " + " | ".join(row) + " |" for row in rows)
                        tables.append({"page": i, "bbox": [], "text": md})
        return _clean_text("\n\n".join(texts)), tables
    except Exception as e:
        logger.exception(f"Failed to extract PPTX {fp}: {e}")
        return "", []


def _df_to_markdown(df, max_rows=500) -> str:
    """DataFrame을 마크다운 테이블로 변환"""
    if len(df) > max_rows:
        df = df.head(max_rows)
    cols = [str(c) for c in df.columns]
    lines = ["| " + " | ".join(cols) + " |"]
    for _, row in df.iterrows():
        lines.append("| " + " | ".join(_clean_text(str(v)) for v in row.tolist()) + " |")
    return "\n".join(lines)


def _extract_csv(fp: Path) -> tuple[str, list[dict]]:
    """CSV 파일 추출"""
    try:
        import pandas as pd
        df = pd.read_csv(fp)
        md = _df_to_markdown(df)
        return "", [{"page": 0, "bbox": [], "text": md}]
    except Exception as e:
        logger.warning(f"Failed to extract CSV {fp}: {e}, trying as plain text")
        return _extract_plain_text(fp)


def _extract_excel(fp: Path) -> tuple[str, list[dict]]:
    """Excel 파일 추출"""
    try:
        import pandas as pd
        xls = pd.ExcelFile(fp)
        tables = []
        for name in xls.sheet_names:
            df = xls.parse(name)
            md = f"### {name}\n" + _df_to_markdown(df)
            tables.append({"page": 0, "bbox": [], "text": md})
        return "", tables
    except Exception as e:
        logger.warning(f"Failed to extract Excel {fp}: {e}")
        return "", []


def _convert_via_libreoffice(src: Path, target_ext: str) -> Optional[Path]:
    """LibreOffice를 통한 문서 변환"""
    try:
        import subprocess
        outdir = src.parent
        subprocess.run([
            "libreoffice", "--headless", "--convert-to", target_ext, 
            "--outdir", str(outdir), str(src)
        ], check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
        cand = src.with_suffix("." + target_ext)
        return cand if cand.exists() else None
    except Exception:
        return None


def _extract_doc(fp: Path) -> tuple[str, list[dict]]:
    """DOC 파일 추출 (LibreOffice 변환 시도)"""
    conv = _convert_via_libreoffice(fp, "docx")
    if conv and conv.exists():
        result = _extract_docx(conv)
        try:
            conv.unlink()  # 임시 변환 파일 삭제
        except Exception:
            pass
        return result
    return _extract_plain_text(fp)


def _extract_ppt(fp: Path) -> tuple[str, list[dict]]:
    """PPT 파일 추출 (LibreOffice 변환 시도)"""
    conv = _convert_via_libreoffice(fp, "pptx")
    if conv and conv.exists():
        result = _extract_pptx(conv)
        try:
            conv.unlink()  # 임시 변환 파일 삭제
        except Exception:
            pass
        return result
    return "", []

def _extract_hwp(fp: Path) -> tuple[str, list[dict]]:
    """HWP 파일 추출 (다단계 폴백)
    순서: python-hwp(API) -> hwp5txt(CLI) -> LibreOffice 변환(docx) -> olefile(구버전 시도)
    반환: (본문텍스트, 표리스트[])
    """
    # --- 0) 공통 헬퍼 ---
    def _try_hwp5txt_cli(path: Path) -> Optional[str]:
        """hwp5txt CLI로 텍스트 추출 (권장 경로)"""
        try:
            import shutil, subprocess, tempfile
            if shutil.which("hwp5txt") is None:
                return None
            # hwp5txt 출력은 stdout. 기본 인코딩은 UTF-8로 가정
            res = subprocess.run(
                ["hwp5txt", str(path)],
                check=True,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                timeout=60,
            )
            out = res.stdout.decode("utf-8", errors="ignore")
            return _clean_text(out) if out else ""
        except Exception:
            return None

    # --- 1) python-hwp 모듈 시도 ---
    try:
        import pyhwp
        from pyhwp.hwp5.xmlmodel import Hwp5File
        texts, tables = [], []
        with Hwp5File(str(fp)) as hwp:
            # 단순 본문 루프 (구조가 달라도 최대한 텍스트 회수)
            for section in getattr(hwp.bodytext, "sections", []):
                # 문단
                for paragraph in getattr(section, "paragraphs", []):
                    try:
                        t = paragraph.get_text()
                        if t and t.strip():
                            texts.append(_clean_text(t))
                    except Exception:
                        continue
                # 표(가능하면)
                for table in getattr(section, "tables", []):
                    try:
                        rows = []
                        for row in getattr(table, "rows", []):
                            cells = []
                            for cell in getattr(row, "cells", []):
                                ctext = cell.get_text() if hasattr(cell, "get_text") else str(cell)
                                cells.append(_clean_text(ctext))
                            if cells:
                                rows.append(cells)
                        if rows:
                            md = "\n".join("| " + " | ".join(r) + " |" for r in rows)
                            tables.append({"page": 0, "bbox": [], "text": md})
                    except Exception:
                        continue
        text_joined = _clean_text("\n\n".join(texts))
        if text_joined or tables:
            return text_joined, tables
    except Exception as e:
        logger.debug(f"python-hwp extraction failed for {fp}: {e}")

    # --- 2) hwp5txt CLI 시도 (가장 잘 되는 편) ---
    cli_text = _try_hwp5txt_cli(fp)
    if cli_text is not None:
        return cli_text, []

    # --- 3) LibreOffice 변환(docx) 시도 (대부분 실패하지만 폴백으로 유지) ---
    conv = _convert_via_libreoffice(fp, "docx")
    if conv and conv.exists():
        try:
            text, tables = _extract_docx(conv)
        finally:
            try:
                conv.unlink()
            except Exception:
                pass
        if text or tables:
            return text, tables

    # --- 4) olefile (구버전 HWP에만 가끔) ---
    try:
        import olefile
        if olefile.isOleFile(str(fp)):
            with olefile.OleFileIO(str(fp)) as ole:
                text_content = ""
                for stream in ole.listdir():
                    try:
                        # 본문 후보 스트림 이름 휴리스틱
                        sname = "/".join(stream)
                        if any(k in sname.lower() for k in ("bodytext", "prvtext", "section", "paragraph")):
                            data = ole.openstream(stream).read()
                            text_content += data.decode("utf-8", errors="ignore")
                    except Exception:
                        continue
                if text_content.strip():
                    return _clean_text(text_content), []
    except Exception as e:
        logger.debug(f"olefile extraction failed for {fp}: {e}")

    logger.warning(f"HWP 파일 추출 실패: {fp}. python-hwp, hwp5txt, LibreOffice, olefile 모두 실패.")
    return "", []


def _extract_any(path: Path) -> tuple[str, list[dict]]:
    """통합 문서 추출 라우터"""
    ext = _ext(path)
    if ext == ".pdf":
        return _extract_pdf_with_tables(path)
    if ext in {".txt", ".text", ".md"}:
        return _extract_plain_text(path)
    if ext == ".docx":
        return _extract_docx(path)
    if ext == ".pptx":
        return _extract_pptx(path)
    if ext == ".csv":
        return _extract_csv(path)
    if ext in {".xlsx", ".xls"}:
        return _extract_excel(path)
    if ext == ".doc":
        return _extract_doc(path)
    if ext == ".ppt":
        return _extract_ppt(path)
    if ext == ".hwp":
        return _extract_hwp(path)
    # 모르는 확장자는 텍스트로 시도
    return _extract_plain_text(path)


# -------------------------------------------------
# 인제스트 파라미터 설정
# -------------------------------------------------
def set_ingest_params(chunk_size: int | None = None, overlap: int | None = None):
    # rag_settings 단일 소스로 저장
    set_vector_settings(chunk_size=chunk_size, overlap=overlap)


def get_ingest_params():
    row = get_vector_settings_row()
    return {"chunkSize": row["chunk_size"], "overlap": row["overlap"]}


# -------------------------------------------------
# Pydantic 스키마
# -------------------------------------------------
class RAGSearchRequest(BaseModel):
    query: str
    top_k: int = Field(5, gt=0)
    user_level: int = Field(1, ge=1)
    task_type: str = Field(..., description="doc_gen | summary | qna")
    model: Optional[str] = None  # 내부적으로 settings에서 로드


class SinglePDFIngestRequest(BaseModel):
    pdf_path: str
    task_types: Optional[List[str]] = None  # 기본은 모든 작업유형
    workspace_id: Optional[int] = None


# -------------------------------------------------
# SQLite 유틸
# -------------------------------------------------


# ====== New helpers ======
def save_raw_file(filename: str, content: bytes) -> str:
    RAW_DATA_DIR.mkdir(parents=True, exist_ok=True)
    out = RAW_DATA_DIR / filename
    out.parent.mkdir(parents=True, exist_ok=True)
    out.write_bytes(content)
    return str(out)


def save_raw_to_row_data(f):
    """Save FastAPI UploadFile to row_data and return relative path."""
    RAW_DATA_DIR.mkdir(parents=True, exist_ok=True)
    name = Path(getattr(f, "filename", "uploaded"))
    dst = RAW_DATA_DIR / name.name
    if dst.exists():
        stem, ext = name.stem, name.suffix
        dst = RAW_DATA_DIR / f"{stem}_{int(time.time())}{ext}"
    with dst.open("wb") as out:
        data = f.file.read() if hasattr(f, "file") else b""
        out.write(data)
    try:
        return str(dst.relative_to(RAW_DATA_DIR))
    except Exception:
        return dst.name


# === Embedding cache(singleton) ===
_EMBED_CACHE: dict[str, tuple[any, any, any]] = {}  # key -> (tok, model, device)
_EMBED_ACTIVE_KEY: Optional[str] = None
_EMBED_LOCK = threading.Lock()

# === Reranker cache(singleton) ===


def _invalidate_embedder_cache():
    global _EMBED_CACHE, _EMBED_ACTIVE_KEY
    with _EMBED_LOCK:
        _EMBED_CACHE.clear()
        _EMBED_ACTIVE_KEY = None


def _get_or_load_embedder(model_key: str, preload: bool = False):
    """
    전역 캐시에서 (tok, model, device) 반환.
    - 캐시에 없으면 로드해서 저장(지연 로딩)
    - preload=True는 의미상 웜업 호출일 뿐, 반환 동작은 동일
    """
    global _EMBED_CACHE, _EMBED_ACTIVE_KEY
    if not model_key:
        raise ValueError(
            "활성화된 임베딩 모델이 없습니다. 먼저 /v1/admin/vector/settings에서 모델을 설정하세요."
        )
    with _EMBED_LOCK:
        if _EMBED_ACTIVE_KEY == model_key and model_key in _EMBED_CACHE:
            return _EMBED_CACHE[model_key]
        _EMBED_CACHE.clear()
        _, model_dir = resolve_model_input(model_key)
        tok, model, device = get_or_load_hf_embedder(str(model_dir))
        _EMBED_CACHE[model_key] = (tok, model, device)
        _EMBED_ACTIVE_KEY = model_key
        return _EMBED_CACHE[model_key]


def warmup_active_embedder(logger_func=print):
    """
    서버 기동 시 호출용(선택). 활성 모델 키를 조회해 캐시를 채움.
    실패해도 서비스는 실제 사용 시 지연 로딩으로 복구됨.
    """
    try:
        key = get_active_embedding_model_name()
        logger_func(f"[warmup] 활성 임베딩 모델: {key}. 로딩 시도...")
        _get_or_load_embedder(key, preload=True)
        logger_func(f"[warmup] 로딩 완료: {key}")
    except Exception as e:
        logger_func(f"[warmup] 로딩 실패(지연 로딩으로 복구 예정): {e}")


async def _get_or_load_embedder_async(model_key: str, preload: bool = False):
    """
    비동기 래퍼: blocking 함수(_get_or_load_embedder)를 스레드풀에서 실행
    이벤트 루프 블로킹 방지
    """
    loop = asyncio.get_running_loop()
    # blocking 함수(_get_or_load_embedder)를 스레드풀에서 실행
    return await loop.run_in_executor(None, _get_or_load_embedder, model_key, preload)



def _set_active_embedding_model(name: str):
    with get_session() as session:
        # 존재하지 않으면 생성
        model = (
            session.query(EmbeddingModel).filter(EmbeddingModel.name == name).first()
        )
        if not model:
            model = EmbeddingModel(name=name, is_active=0)
            session.add(model)
            session.flush()
        # 모두 비활성 → 대상만 활성
        session.query(EmbeddingModel).filter(EmbeddingModel.is_active == 1).update(
            {"is_active": 0, "activated_at": None}
        )
        model.is_active = 1
        model.activated_at = now_kst()
        session.commit()


def _update_vector_settings(
    search_type: Optional[str] = None,
    chunk_size: Optional[int] = None,
    overlap: Optional[int] = None,
):
    """레거시 API 호환: rag_settings(싱글톤) 업데이트"""
    cur = get_vector_settings_row()
    new_search = (search_type or cur["search_type"]).lower()
    if new_search == "vector":
        new_search = "semantic"
    if new_search not in {"hybrid", "semantic", "bm25"}:
        raise ValueError(
            "unsupported searchType; allowed: 'hybrid','semantic','bm25' (or 'vector' alias)"
        )
    new_chunk = int(chunk_size if chunk_size is not None else cur["chunk_size"])
    new_overlap = int(overlap if overlap is not None else cur["overlap"])
    if new_chunk <= 0 or new_overlap < 0 or new_overlap >= new_chunk:
        raise ValueError("invalid chunk/overlap (chunk>0, 0 <= overlap < chunk)")

    with get_session() as session:
        s = session.query(RagSettings).filter(RagSettings.id == 1).first()
        if not s:
            s = RagSettings(id=1)
            session.add(s)
        s.search_type = new_search
        s.chunk_size = new_chunk
        s.overlap = new_overlap
        s.updated_at = now_kst()
        session.commit()


def _milvus_has_data(collection_name: str = COLLECTION_NAME) -> bool:
    client = _client()
    if collection_name not in client.list_collections():
        return False
    try:
        rows = client.query(collection_name=collection_name, output_fields=["pk"], limit=1)
        return len(rows) > 0
    except Exception:
        return True


# ---------------- Vector Settings ----------------
def set_vector_settings(embed_model_key: Optional[str] = None,
                        search_type: Optional[str] = None,
                        chunk_size: Optional[int] = None,
                        overlap: Optional[int] = None) -> Dict:
    """
    rag_settings 단일 소스로 설정 저장.
    - 임베딩 모델 변경 시 기존 데이터 존재하면 차단, 활성 모델 갱신 및 캐시 무효화
    - search_type/청크/오버랩은 rag_settings에만 반영
    """
    cur = get_vector_settings()
    key_now = cur.get("embeddingModel")
    st_now = (cur.get("searchType") or "hybrid").lower()
    cs_now = int(cur.get("chunkSize") or 512)
    ov_now = int(cur.get("overlap") or 64)

    new_key = embed_model_key or key_now
    new_st = (search_type or st_now).lower()
    # DB 제약과 일치(semantic == vector)
    if new_st == "semantic":
        new_st = "vector"
    if new_st not in {"hybrid", "bm25", "vector"}:
        raise ValueError("unsupported searchType; allowed: 'hybrid','bm25','vector'")

    new_cs = int(chunk_size if chunk_size is not None else cs_now)
    new_ov = int(overlap if overlap is not None else ov_now)
    if new_cs <= 0 or new_ov < 0 or new_ov >= new_cs:
        raise ValueError("invalid chunk/overlap (chunk>0, 0 <= overlap < chunk)")

    if embed_model_key is not None:
        if _milvus_has_data():
            raise RuntimeError("Milvus 컬렉션에 기존 데이터가 남아있습니다. 먼저 /v1/admin/vector/delete-all 을 호출해 초기화하세요.")
        _set_active_embedding_model(embed_model_key)
        _invalidate_embedder_cache()

    with get_session() as session:
        s = session.query(RagSettings).filter(RagSettings.id == 1).first()
        if not s:
            s = RagSettings(id=1)
            session.add(s)
        s.embedding_key = new_key
        # search_type/chunk/overlap은 _update_vector_settings에서 반영됨. 여기선 존재 시 보존
        if search_type is not None:
            s.search_type = (
                (search_type or "hybrid").lower().replace("vector", "semantic")
            )
        if chunk_size is not None:
            s.chunk_size = int(chunk_size)
        if overlap is not None:
            s.overlap = int(overlap)
        s.updated_at = now_kst()
        session.commit()

    return get_vector_settings()


def get_vector_settings() -> Dict:
    # rag_settings 는 검색 타입/청크/오버랩만 신뢰
    row = get_vector_settings_row()  # {"search_type": "...", "chunk_size": 512, "overlap": 64}
    try:
        # ★활성 모델만 신뢰 (EmbeddingModel.is_active == 1)
        model = get_active_embedding_model_name()
    except Exception:
        model = None

    return {
        "embeddingModel": model,                        # ← rag_settings.embedding_key는 무시
        "searchType": row.get("search_type", "hybrid"),
        "chunkSize": int(row.get("chunk_size", 512)),
        "overlap": int(row.get("overlap", 64)),
    }


def list_available_embedding_models() -> List[str]:
    """
    ./storage/embedding-models 폴더 내의 모델 폴더명들을 반환.
    - embedding_ 접두사가 있으면 제거 (예: embedding_bge_m3 → bge_m3)
    - 폴더만 반환 (파일 제외)
    """
    models = []
    if not MODEL_ROOT_DIR.exists():
        return models
    
    for item in MODEL_ROOT_DIR.iterdir():
        if item.is_dir():
            model_name = item.name
            # embedding_ 접두사 제거
            if model_name.startswith("embedding_"):
                model_name = model_name[len("embedding_"):]
            models.append(model_name)
    
    return sorted(models)

# ------------- Security Level (per task) ---------
def _parse_at_string_to_keywords(value: str) -> List[str]:
    if not value:
        return []
    toks = [t.strip() for t in value.split("@")]
    return [t for t in toks if t]


def _normalize_keywords(val: Any) -> List[str]:
    """
    리스트/튜플/셋: 각 원소를 str로 캐스팅, 공백/해시 제거
    문자열: '@' 기준으로 토큰화
    빈 값 제거 및 중복 제거
    """
    out: List[str] = []
    if isinstance(val, str):
        toks = [t.strip() for t in val.split("@")]
    elif isinstance(val, (list, tuple, set)):
        toks = [str(t).strip() for t in val]
    else:
        toks = []
    for t in toks:
        if not t:
            continue
        if t.startswith("#"):
            t = t[1:]
        if t and t not in out:
            out.append(t)
    return out


def _normalize_levels(
    levels_raw: Dict[str, Any], max_level: int
) -> Dict[int, List[str]]:
    norm: Dict[int, List[str]] = {}
    for k, v in (levels_raw or {}).items():
        try:
            lv = int(str(k).strip().replace("level_", ""))
        except Exception:
            continue
        if lv < 1 or lv > max_level:
            continue
        kws = _normalize_keywords(v)
        if kws:
            norm[lv] = kws
    return norm


def upsert_security_level_for_task(
    task_type: str, max_level: int, levels_raw: Dict[str, Any]
) -> Dict:
    if task_type not in TASK_TYPES:
        raise ValueError(f"invalid task_type: {task_type}")
    if max_level < 1:
        raise ValueError("maxLevel must be >= 1")

    levels_map = _normalize_levels(levels_raw, max_level)

    with get_session() as session:
        # upsert config
        cfg = (
            session.query(SecurityLevelConfigTask)
            .filter(SecurityLevelConfigTask.task_type == task_type)
            .first()
        )
        if not cfg:
            cfg = SecurityLevelConfigTask(task_type=task_type, max_level=int(max_level))
            session.add(cfg)
        else:
            cfg.max_level = int(max_level)
            cfg.updated_at = now_kst()
        # replace keywords
        session.query(SecurityLevelKeywordsTask).filter(
            SecurityLevelKeywordsTask.task_type == task_type
        ).delete()
        for lv, kws in levels_map.items():
            for kw in kws:
                session.add(
                    SecurityLevelKeywordsTask(
                        task_type=task_type, level=int(lv), keyword=str(kw)
                    )
                )
        session.commit()
        return get_security_level_rules_for_task(task_type)


def get_security_level_rules_for_task(task_type: str) -> Dict:
    with get_session() as session:
        cfg = (
            session.query(SecurityLevelConfigTask)
            .filter(SecurityLevelConfigTask.task_type == task_type)
            .first()
        )
        max_level = int(cfg.max_level) if cfg else 1
        res: Dict[str, Any] = {
            "taskType": task_type,
            "maxLevel": max_level,
            "levels": {str(i): [] for i in range(1, max_level + 1)},
        }
        rows = (
            session.query(
                SecurityLevelKeywordsTask.level, SecurityLevelKeywordsTask.keyword
            )
            .filter(SecurityLevelKeywordsTask.task_type == task_type)
            .order_by(
                SecurityLevelKeywordsTask.level.asc(),
                SecurityLevelKeywordsTask.keyword.asc(),
            )
            .all()
        )
        for lv, kw in rows:
            key = str(int(lv))
            res["levels"].setdefault(key, []).append(str(kw))
        return res


def set_security_level_rules_per_task(config: Dict[str, Dict]) -> Dict:
    """
    config = {
      "doc_gen": {"maxLevel": 3, "levels": {"2": "@금액@연봉", "3": "@부정@퇴직금"}},
      "summary": {"maxLevel": 2, "levels": {"2": "@사내비밀"}},
      "qna": {"maxLevel": 3, "levels": {"2": "@연구", "3": "@개인정보"}}
    }
    """
    with get_session() as session:
        # 전체 삭제 후 재삽입(간결/명확)
        session.query(SecurityLevelConfigTask).delete()
        session.query(SecurityLevelKeywordsTask).delete()
        session.flush()

        for task in TASK_TYPES:
            entry = config.get(task) or {}
            max_level = int(entry.get("maxLevel", 1))
            session.add(
                SecurityLevelConfigTask(task_type=task, max_level=max(1, max_level))
            )
            levels = entry.get("levels", {}) or {}
            for lvl_str, at_str in levels.items():
                try:
                    lvl = int(str(lvl_str).strip().replace("level_", ""))
                except Exception:
                    continue
                if lvl <= 1 or lvl > max_level:
                    continue
                for kw in _parse_at_string_to_keywords(str(at_str)):
                    session.add(
                        SecurityLevelKeywordsTask(
                            task_type=task, level=int(lvl), keyword=str(kw)
                        )
                    )
        session.commit()
        return get_security_level_rules_all()


def get_security_level_rules_all() -> Dict:
    with get_session() as session:
        # 기본 max_level=1
        max_map = {t: 1 for t in TASK_TYPES}
        for task, max_level in session.query(
            SecurityLevelConfigTask.task_type, SecurityLevelConfigTask.max_level
        ).all():
            max_map[task] = int(max_level)

        res: Dict[str, Dict] = {}
        for task in TASK_TYPES:
            res[task] = {
                "maxLevel": max_map.get(task, 1),
                "levels": {str(i): [] for i in range(1, max_map.get(task, 1) + 1)},
            }

        rows = (
            session.query(
                SecurityLevelKeywordsTask.task_type,
                SecurityLevelKeywordsTask.level,
                SecurityLevelKeywordsTask.keyword,
            )
            .order_by(
                SecurityLevelKeywordsTask.task_type.asc(),
                SecurityLevelKeywordsTask.level.asc(),
                SecurityLevelKeywordsTask.keyword.asc(),
            )
            .all()
        )
        for task, level, kw in rows:
            if task in res:
                lv = str(int(level))
                if lv not in res[task]["levels"]:
                    res[task]["levels"][lv] = []
                res[task]["levels"][lv].append(str(kw))
        return res


def _determine_level_for_task(text: str, task_rules: Dict) -> int:
    max_level = int(task_rules.get("maxLevel", 1))
    levels = task_rules.get("levels", {})
    sel = 1
    # 상위 레벨 우선
    for lvl in range(1, max_level + 1):
        kws = levels.get(str(lvl), [])
        for kw in kws:
            if kw and kw in text:
                sel = max(sel, lvl)
    return sel


# --- add: test 컬렉션 전용 보조 함수 ---
def _ensure_collection_and_index_for(
    client: MilvusClient,
    collection_name: str,
    emb_dim: int,
    metric: str = "IP",
):
    """
    _ensure_collection_and_index의 'collection_name' 파라미터 버전 (세션 컬렉션용)
    """
    logger.info(f"[Milvus] 컬렉션/인덱스 준비: {collection_name}")

    cols = client.list_collections()
    if collection_name not in cols:
        logger.info(f"[Milvus] 컬렉션 생성: {collection_name}")
        schema = client.create_schema(
            auto_id=True, enable_dynamic_field=False, description=f"PDF chunks ({collection_name})"
        )
        schema.add_field("pk", DataType.INT64, is_primary=True)
        schema.add_field("embedding", DataType.FLOAT_VECTOR, dim=int(emb_dim))
        schema.add_field("path", DataType.VARCHAR, max_length=500)
        schema.add_field("chunk_idx", DataType.INT64)
        schema.add_field("task_type", DataType.VARCHAR, max_length=16)  # doc_gen|summary|qna
        schema.add_field("security_level", DataType.INT64)
        schema.add_field("doc_id", DataType.VARCHAR, max_length=255)
        schema.add_field("version", DataType.INT64)
        schema.add_field("page", DataType.INT64)  # 페이지 번호

        # text / text_sparse (BM25)
        try:
            schema.add_field("text", DataType.VARCHAR, max_length=32768, enable_analyzer=True)
        except TypeError:
            schema.add_field("text", DataType.VARCHAR, max_length=32768)
        try:
            schema.add_field("text_sparse", DataType.SPARSE_FLOAT_VECTOR)
        except Exception:
            logger.warning("[Milvus] SPARSE_FLOAT_VECTOR 미지원 클라이언트 - 서버 BM25 하이브리드 불가")

        # BM25 Function
        if Function is not None:
            try:
                fn = Function(
                    name="bm25_text2sparse",
                    function_type=FunctionType.BM25,
                    input_field_names=["text"],
                    output_field_names=["text_sparse"],
                )
                schema.add_function(fn)
                logger.info("[Milvus] BM25 Function 연결 완료 (text -> text_sparse)")
            except Exception as e:
                logger.warning(f"[Milvus] BM25 Function 추가 실패: {e}")

        client.create_collection(collection_name=collection_name, schema=schema)
        logger.info(f"[Milvus] 컬렉션 생성 완료: {collection_name}")

    # 1) dense index
    try:
        idx_dense = client.list_indexes(collection_name=collection_name, field_name="embedding")
    except Exception:
        idx_dense = []
    if not idx_dense:
        ip = client.prepare_index_params()
        ip.add_index("embedding", "FLAT", metric_type=metric, params={})
        client.create_index(collection_name, ip, timeout=180.0, sync=True)

    # 2) sparse index
    try:
        idx_sparse = client.list_indexes(collection_name=collection_name, field_name="text_sparse")
    except Exception:
        idx_sparse = []
    if not idx_sparse:
        ip2 = client.prepare_index_params()
        try:
            ip2.add_index("text_sparse", "SPARSE_INVERTED_INDEX", params={})
        except TypeError:
            ip2.add_index("text_sparse", "SPARSE_INVERTED_INDEX", metric_type="BM25", params={})
        client.create_index(collection_name, ip2, timeout=180.0, sync=True)

    # reload
    try:
        client.release_collection(collection_name=collection_name)
    except Exception:
        pass
    client.load_collection(collection_name=collection_name)
    logger.info(f"[Milvus] 로드 완료: {collection_name}")


# -------------------------------------------------
# Milvus Client / 컬렉션 스키마
# -------------------------------------------------
def _client() -> MilvusClient:
    kwargs = {"uri": MILVUS_URI}
    if MILVUS_TOKEN:
        kwargs["token"] = MILVUS_TOKEN
    return MilvusClient(**kwargs)
def _ensure_collection_and_index(
    client: MilvusClient,
    emb_dim: int,
    metric: str = "IP",
    collection_name: str = COLLECTION_NAME,
):
    logger.info(f"[Milvus] 컬렉션 및 인덱스 준비 시작: {collection_name}")
    cols = client.list_collections()
    if collection_name not in cols:
        logger.info(f"[Milvus] 컬렉션 생성: {collection_name}")
        schema = client.create_schema(
            auto_id=True, enable_dynamic_field=False, description=f"PDF chunks ({collection_name})"
        )
        schema.add_field("pk", DataType.INT64, is_primary=True)
        schema.add_field("embedding", DataType.FLOAT_VECTOR, dim=int(emb_dim))
        schema.add_field("path", DataType.VARCHAR, max_length=500)
        schema.add_field("chunk_idx", DataType.INT64)
        schema.add_field("task_type", DataType.VARCHAR, max_length=16)  # 'doc_gen'|'summary'|'qna'
        schema.add_field("security_level", DataType.INT64)
        schema.add_field("doc_id", DataType.VARCHAR, max_length=255)
        schema.add_field("version", DataType.INT64)
        schema.add_field("page", DataType.INT64)  # 페이지 번호
        # 하이브리드용 텍스트/스파스 필드
        try:
            schema.add_field("text", DataType.VARCHAR, max_length=32768, enable_analyzer=True)
        except TypeError:
            schema.add_field("text", DataType.VARCHAR, max_length=32768)
        try:
            schema.add_field("text_sparse", DataType.SPARSE_FLOAT_VECTOR)
        except Exception:
            logger.warning("[Milvus] SPARSE_FLOAT_VECTOR 미지원 클라이언트입니다. 서버 BM25 하이브리드 사용 불가.")

        if Function is not None:
            try:
                fn = Function(
                    name="bm25_text2sparse",
                    function_type=FunctionType.BM25,
                    input_field_names=["text"],
                    output_field_names=["text_sparse"],
                )
                schema.add_function(fn)
                logger.info("[Milvus] BM25 Function 연결 완료 (text -> text_sparse)")
            except Exception as e:
                logger.warning(f"[Milvus] BM25 Function 추가 실패: {e}")
        client.create_collection(collection_name=collection_name, schema=schema)
        logger.info(f"[Milvus] 컬렉션 생성 완료: {collection_name}")

    # 1) 덴스 벡터 인덱스
    try:
        idx_dense = client.list_indexes(collection_name=collection_name, field_name="embedding")
    except Exception:
        idx_dense = []
    if not idx_dense:
        logger.info(f"[Milvus] (embedding) 인덱스 생성 시작 @ {collection_name}")
        ip = client.prepare_index_params()
        ip.add_index("embedding", "FLAT", metric_type=metric, params={})
        client.create_index(collection_name, ip, timeout=180.0, sync=True)
        logger.info(f"[Milvus] (embedding) 인덱스 생성 완료 @ {collection_name}")

    # 2) 스파스 인덱스
    try:
        idx_sparse = client.list_indexes(collection_name=collection_name, field_name="text_sparse")
    except Exception:
        idx_sparse = []
    if not idx_sparse:
        logger.info(f"[Milvus] (text_sparse) 인덱스 생성 시작 @ {collection_name}")
        ip2 = client.prepare_index_params()
        try:
            ip2.add_index("text_sparse", "SPARSE_INVERTED_INDEX", params={})
        except TypeError:
            ip2.add_index("text_sparse", "SPARSE_INVERTED_INDEX", metric_type="BM25", params={})
        client.create_index(collection_name, ip2, timeout=180.0, sync=True)
        logger.info(f"[Milvus] (text_sparse) 인덱스 생성 완료 @ {collection_name}")

    # 로드
    try:
        client.release_collection(collection_name=collection_name)
    except Exception:
        pass
    client.load_collection(collection_name=collection_name)
    logger.info(f"[Milvus] 컬렉션 로드 완료: {collection_name}")


# -------------------------------------------------
# 1) PDF → 텍스트 추출 (작업유형별 보안레벨 동시 산정)
# -------------------------------------------------
async def extract_pdfs():
    from tqdm import tqdm  # type: ignore

    EXTRACTED_TEXT_DIR.mkdir(parents=True, exist_ok=True)
    LOCAL_DATA_ROOT.mkdir(parents=True, exist_ok=True)
    RAW_DATA_DIR.mkdir(parents=True, exist_ok=True)

    # 규칙 로드
    all_rules = get_security_level_rules_all()  # {task: {"maxLevel":N, "levels":{...}}}

    # 이전 메타 로드
    prev_meta: Dict[str, Dict] = {}
    if META_JSON_PATH.exists():
        try:
            prev_meta = json.loads(META_JSON_PATH.read_text(encoding="utf-8"))
        except Exception:
            logger.exception("Failed to read META JSON; recreating.")

    # 중복 제거 로직: 파일명 마지막 토큰이 날짜/버전 숫자면 최신만 유지
    def _extract_base_and_date(p: Path):
        name = p.stem
        parts = name.split("_")
        date_num = 0
        if len(parts) >= 2:
            cand = parts[-1]
            if cand.isdigit() and len(cand) in (4, 6, 8):
                try:
                    date_num = int(cand)
                except Exception:
                    date_num = 0
        mid_tokens = [t for t in parts[:-1] if t and not t.isdigit()]
        base = max(mid_tokens, key=len) if mid_tokens else parts[0]
        return base, date_num

    raw_files = [p for p in RAW_DATA_DIR.rglob("*") if p.is_file() and _ext(p) in SUPPORTED_EXTS]

    # base(문서ID 유사)별로 버전 후보 묶기: (Path, date_num)
    grouped: Dict[str, List[Tuple[Path, int]]] = defaultdict(list)
    for p in raw_files:
        base, date_num = _extract_base_and_date(p)
        grouped[base].append((p, date_num))

    kept, removed = [], []
    for base, lst in grouped.items():
        # lst 원소는 (Path, date_num)
        lst_sorted = sorted(
            lst,
            key=lambda it: (it[1], it[0].stat().st_mtime, len(it[0].name))  # date_num → mtime → 이름 길이
        )
        keep_path = lst_sorted[-1][0]          # 최신 후보의 Path
        kept.append(keep_path)
        for old_path, _ in lst_sorted[:-1]:
            try:
                old_path.unlink(missing_ok=True)
                removed.append(str(old_path.relative_to(RAW_DATA_DIR)))
            except Exception:
                logger.exception("Failed to remove duplicate: %s", old_path)

    if not kept:
        return {
            "message": "처리할 문서가 없습니다.",
            "meta_path": str(META_JSON_PATH),
            "deduplicated": {"removedCount": len(removed), "removed": removed},
        }

    new_meta: Dict[str, Dict] = {}
    for src in tqdm(kept, desc="문서 전처리"):
        try:
            logger.info(f"[Extract] 파일 처리 시작: {src.name}")
            
            # PDF인 경우 페이지별 정보를 포함하여 추출
            pages_text_dict: Dict[int, str] = {}
            total_pages = 0
            if _ext(src) == ".pdf":
                text, tables, pages_text_dict, total_pages = _extract_pdf_with_tables(src)
            else:
                text, tables = _extract_any(src)
            
            # 표 추출 결과 로깅
            if tables:
                logger.info(f"[Extract] {src.name}: 표 {len(tables)}개 추출됨")
                for t_idx, t in enumerate(tables):
                    page = t.get("page", 0)
                    text_preview = (t.get("text", "")[:100] + "...") if t.get("text") else ""
                    logger.info(f"[Extract] {src.name} 표 {t_idx+1}: 페이지={page}, 텍스트 미리보기={text_preview}")
            else:
                logger.info(f"[Extract] {src.name}: 추출된 표 없음 (텍스트만 추출됨)")
            
            # 작업유형별 보안 레벨 (본문+표 모두 포함해서 판정)
            whole_for_level = text + "\n\n" + "\n\n".join(t.get("text","") for t in (tables or []))
            sec_map = {task: _determine_level_for_task(
                whole_for_level, all_rules.get(task, {"maxLevel": 1, "levels": {}})
            ) for task in TASK_TYPES}

            max_sec = max(sec_map.values()) if sec_map else 1
            sec_folder = f"securityLevel{int(max_sec)}"

            rel_from_raw = src.relative_to(RAW_DATA_DIR)
            # 원본 그대로 보관
            dest_rel = Path(sec_folder) / rel_from_raw
            dest_abs = LOCAL_DATA_ROOT / dest_rel
            dest_abs.parent.mkdir(parents=True, exist_ok=True)
            try:
                shutil.copy2(src, dest_abs)
            except Exception:
                logger.exception("Failed to copy file: %s", dest_abs)

            # 통합 파일 저장 (.txt, 마크다운 형식) - 페이지 순서대로 텍스트와 표 배치
            # 파일명은 원본 파일명 그대로 사용 (확장자만 .txt)
            txt_rel = dest_rel.with_suffix(".txt")
            (EXTRACTED_TEXT_DIR / txt_rel).parent.mkdir(parents=True, exist_ok=True)
            saved_files: Dict[str, str] = {}
            
            # 통합 파일 저장
            combined_txt_file = EXTRACTED_TEXT_DIR / txt_rel
            try:
                # 페이지별 표 그룹화
                pages_tables: Dict[int, list[dict]] = defaultdict(list)
                for t in (tables or []):
                    page_num = t.get("page", 0)
                    if page_num > 0:
                        pages_tables[page_num].append(t)
                
                # 통합 파일 작성 (페이지 순서대로) - 페이지 마커 없이 순수 텍스트+표만 저장
                with open(combined_txt_file, "w", encoding="utf-8") as f:
                    # PDF인 경우: 페이지별 텍스트 정보 사용
                    if pages_text_dict:
                        # 모든 페이지 번호 수집 (텍스트와 표 모두 고려)
                        all_page_nums = set(pages_text_dict.keys())
                        all_page_nums.update(pages_tables.keys())
                        
                        if all_page_nums:
                            for page_num in sorted(all_page_nums):
                                # 해당 페이지의 텍스트 (페이지 마커 없이)
                                page_text_content = pages_text_dict.get(page_num, "")
                                if page_text_content:
                                    f.write(page_text_content)
                                    f.write("\n\n")
                                
                                # 해당 페이지의 표들 (텍스트 뒤에 삽입)
                                page_tables_list = pages_tables.get(page_num, [])
                                if page_tables_list:
                                    for t_idx, t in enumerate(page_tables_list):
                                        table_text = t.get("text", "")
                                        if table_text:
                                            f.write(table_text)
                                            f.write("\n\n")
                                
                                # 페이지 구분선 (마지막 페이지가 아니면)
                                if page_num < max(all_page_nums):
                                    f.write("\n---\n\n")
                        else:
                            # 페이지 정보가 없으면 전체 텍스트만
                            if text.strip():
                                f.write(text)
                                f.write("\n\n")
                            if tables:
                                for t in tables:
                                    table_text = t.get("text", "")
                                    if table_text:
                                        f.write(table_text)
                                        f.write("\n\n")
                    else:
                        # PDF가 아닌 경우: 전체 텍스트와 표만 저장
                        if text.strip():
                            f.write(text)
                            f.write("\n\n")
                        if tables:
                            for t in tables:
                                table_text = t.get("text", "")
                                if table_text:
                                    f.write(table_text)
                                    f.write("\n\n")
                
                saved_files["text"] = str(combined_txt_file)
                logger.info(f"[Extract] 통합 파일 저장: {combined_txt_file}")
            except Exception as e:
                logger.exception(f"[Extract] 통합 파일 저장 실패: {e}")

            # doc_id/version 유추
            stem = rel_from_raw.stem
            doc_id, version = _parse_doc_version(stem)

            info = {
                "chars": len(text),
                "lines": len(text.splitlines()),
                "preview": (_clean_text(text[:200].replace("\n"," ")) + "…") if text else "",
                "security_levels": sec_map,  # 작업유형별 보안레벨
                "doc_id": doc_id,
                "version": version,
                "tables": tables or [],  # ★ 표 정보 추가
                "pages": pages_text_dict if pages_text_dict else {},  # ★ 페이지별 텍스트 정보 (메타데이터용)
                "total_pages": total_pages,  # ★ 총 페이지 수
                "sourceExt": _ext(src),  # 원본 확장자 기록
                "saved_files": saved_files,  # 저장된 파일 경로 추가
                # 페이로드 정보 (LLM에 전달하지 않지만 메타데이터로 저장)
                "extraction_info": {
                    "original_file": src.name,
                    "text_length": len(text),
                    "table_count": len(tables or []),
                    "extracted_at": now_kst_string(),
                },
            }
            new_meta[str(dest_rel)] = info
            logger.info(f"[Extract] {src.name}: 메타데이터 저장 완료 (텍스트={len(text)}자, 표={len(tables or [])}개)")

        except Exception as e:
            logger.exception("Failed to process: %s", src)
            try:
                rel_from_raw = src.relative_to(RAW_DATA_DIR)
                dest_rel = Path("securityLevel1") / rel_from_raw
                new_meta[str(dest_rel)] = {"error": str(e)}
            except Exception:
                new_meta[src.name] = {"error": str(e)}

    META_JSON_PATH.write_text(
        json.dumps(new_meta, ensure_ascii=False, indent=2), encoding="utf-8"
    )
    return {
        "message": "문서 추출 완료",
        "file_count": len(kept),
        "meta_path": str(META_JSON_PATH),
        "deduplicated": {"removedCount": len(removed), "removed": removed},
    }


def _parse_doc_version(stem: str) -> Tuple[str, int]:
    if "_" in stem:
        base, cand = stem.rsplit("_", 1)
        if cand.isdigit() and len(cand) in (4, 8):
            return base, int(cand)
    return stem, 0


# -------------------------------------------------
# 2) 인제스트 (bulk)
#   - 작업유형별로 동일 청크를 각각 저장(task_type, security_level 분리)
# -------------------------------------------------
async def ingest_embeddings(
    model_key: str | None = None,
    chunk_size: int | None = None,
    overlap: int | None = None,
    target_tasks: list[str] | None = None,
    collection_name: str = COLLECTION_NAME,
    file_keys_filter: list[str] | None = None,  # ★ 추가: 특정 파일만 인제스트
):
    """
    META_JSON을 읽어 추출된 텍스트(.txt)들을 인제스트한다.
    - VARCHAR(32768 bytes) 초과 방지: _split_for_varchar_bytes 로 안전 분할
    - 표는 [[TABLE ...]] 머리글 유지, 이어지는 조각은 [[TABLE_CONT i/n]] 마커로 연속성 표시
    - collection_name 파라미터를 끝까지 사용(기본/세션 컬렉션 공용)
    - file_keys_filter 가 주어지면 해당되는 파일(meta key/파일명/스텀)이 '포함'된 항목만 인제스트
    """
    # ==== 설정/모델 ====
    settings = get_vector_settings()
    MAX_TOKENS = int(chunk_size if chunk_size is not None else settings["chunkSize"])
    OVERLAP = int(overlap if overlap is not None else settings["overlap"])

    if not META_JSON_PATH.exists():
        return {"error": "메타 JSON이 없습니다. 먼저 PDF/문서 추출을 수행하세요."}

    eff_model_key = model_key or settings["embeddingModel"]
    tok, model, device = await _get_or_load_embedder_async(eff_model_key)
    
    # 벡터 차원 검증
    probe_vec = hf_embed_text(tok, model, device, "probe")
    emb_dim = int(probe_vec.shape[0])
    logger.info(f"[Ingest] 임베딩 모델: {eff_model_key}, 벡터 차원: {emb_dim}")
    
    client = _client()
    
    # 기존 컬렉션이 있으면 차원을 확인하고, 다르면 삭제
    if collection_name in client.list_collections():
        try:
            # 컬렉션 정보 확인
            desc = client.describe_collection(collection_name)
            existing_dim = None
            for field in desc.get("fields", []):
                if field.get("name") == "embedding":
                    existing_dim = field.get("params", {}).get("dim")
                    break
            
            if existing_dim and int(existing_dim) != emb_dim:
                logger.warning(f"[Ingest] 차원 불일치: 기존={existing_dim}, 새모델={emb_dim}. 컬렉션 재생성.")
                client.drop_collection(collection_name)
        except Exception as e:
            logger.warning(f"[Ingest] 컬렉션 정보 확인 실패: {e}. 재생성 시도.")
            try:
                client.drop_collection(collection_name)
            except Exception:
                pass
    
    _ensure_collection_and_index(client, emb_dim, metric="IP", collection_name=collection_name)

    # ==== META 로드 및 대상 필터 구성 ====
    meta: dict = json.loads(META_JSON_PATH.read_text(encoding="utf-8"))
    tasks = [t for t in (target_tasks or TASK_TYPES) if t in TASK_TYPES]
    if not tasks:
        return {"error": f"유효한 작업유형이 없습니다. 허용: {TASK_TYPES}"}

    filter_tokens = set()
    if file_keys_filter:
        # meta key / 파일명 / 스템을 모두 매칭할 수 있도록 소문자 토큰화
        for f in file_keys_filter:
            p = Path(str(f))
            filter_tokens.add(str(f).lower())
            filter_tokens.add(p.name.lower())
            filter_tokens.add(p.stem.lower())

    total_inserted = 0
    BATCH_SIZE = 128

    # ==== 인제스트 ====
    # 주의: EXTRACTED_TEXT_DIR 안의 *.txt 를 돌면서, 해당 txt 가 어떤 meta key(원본 확장자)와 매칭되는지 찾는다.
    for txt_path in EXTRACTED_TEXT_DIR.rglob("*.txt"):
        rel_txt = txt_path.relative_to(EXTRACTED_TEXT_DIR)

        # 다양한 확장자 후보로 META key 찾기
        cands = [rel_txt.with_suffix(ext).as_posix() for ext in SUPPORTED_EXTS]
        meta_key = next((k for k in cands if k in meta), None)
        if not meta_key:
            continue

        # ★ 업로드한 것만 인제스트 옵션: meta key / 파일명 / 스템 기준 필터링
        if filter_tokens:
            p = Path(meta_key)
            if (meta_key.lower() not in filter_tokens and
                p.name.lower() not in filter_tokens and
                p.stem.lower() not in filter_tokens):
                continue

        entry = meta.get(meta_key) or {}
        sec_map = entry.get("security_levels", {}) or {}

        # doc_id / version 확보(없으면 파일명에서 유추)
        doc_id = entry.get("doc_id")
        version = int(entry.get("version", 0) or 0)
        if not doc_id or version == 0:
            _id, _ver = _parse_doc_version(Path(meta_key).stem)
            doc_id = doc_id or _id
            version = version or _ver
            entry["doc_id"] = doc_id
            entry["version"] = version
            meta[meta_key] = entry  # 변경사항 반영

        # 기존 동일 문서/버전 삭제(작업유형 상관 없이)
        try:
            client.delete(
                collection_name=collection_name,
                filter=f"doc_id == '{doc_id}' && version <= {int(version)}",
            )
        except Exception:
            pass

        # 본문 텍스트 로드 및 청크화
        try:
            text = txt_path.read_text(encoding="utf-8")
        except Exception:
            # 혹시 모를 인코딩 문제 폴백
            text = txt_path.read_text(errors="ignore")
        
        # 통합 파일을 직접 파싱하여 페이지별로 분할 (텍스트와 표가 함께 저장된 파일)
        # 페이지 구분선 "---" 기준으로 페이지 분리
        def _parse_integrated_file(text: str) -> list[tuple[int, str]]:
            """통합 파일을 페이지별로 분할 (페이지 구분선 "---" 기준)"""
            page_blocks: list[tuple[int, str]] = []
            lines = text.split('\n')
            current_page = 1
            current_content = []
            
            for line in lines:
                # 페이지 구분선 확인: "---" (빈 줄로 둘러싸인 경우)
                if line.strip() == "---":
                    # 이전 페이지 저장
                    if current_content:
                        page_text = '\n'.join(current_content).strip()
                        if page_text:
                            page_blocks.append((current_page, page_text))
                    current_page += 1
                    current_content = []
                else:
                    current_content.append(line)
            
            # 마지막 페이지 저장
            if current_content:
                page_text = '\n'.join(current_content).strip()
                if page_text:
                    page_blocks.append((current_page, page_text))
            
            # 페이지 구분선이 없으면 전체를 1페이지로 처리
            if not page_blocks:
                if text.strip():
                    page_blocks = [(1, text.strip())]
            
            return page_blocks
        
        # 통합 파일 파싱
        page_blocks = _parse_integrated_file(text)
        logger.info(f"[Ingest] 통합 파일 파싱: {len(page_blocks)}개 페이지 블록 발견")
        
        # 전체 문서에서 청크 인덱스 누적 (페이지별로 0부터 시작하지 않도록)
        chunks_with_page: list[tuple[int, int, str]] = []  # (page, chunk_idx, chunk_text)
        global_chunk_idx = 0  # 전체 문서에서 누적되는 청크 인덱스
        
        for page_num, page_text in page_blocks:
            if not page_text:
                continue
            page_chunks = chunk_text(page_text)
            for chunk in page_chunks:
                if chunk.strip():  # 빈 청크 제외
                    chunks_with_page.append((page_num, global_chunk_idx, chunk))
                    global_chunk_idx += 1
        
        logger.info(f"[Ingest] 총 {global_chunk_idx}개 청크 생성 (페이지별 청크 인덱스 누적)")
        
        # 표 블록 처리
        # 통합 파일에 이미 표가 포함되어 있으므로, 표를 별도로 인제스트하지 않음
        # (통합 파일을 파싱할 때 표도 함께 청크화되므로 중복 방지)
        tables = entry.get("tables", []) or []
        logger.info(f"[Ingest] 표 정보: {len(tables)}개 (통합 파일에 이미 포함되어 있으므로 별도 인제스트 안 함)")

        batch: list[dict] = []

        for task in tasks:
            lvl = int(sec_map.get(task, 1))

            # 1) 본문 조각 (페이지 정보 포함, 텍스트와 표 모두 포함)
            for page_num, idx, c in chunks_with_page:
                # VARCHAR 한도 안전 분할(바이트 기준)
                for part in _split_for_varchar_bytes(c):
                    # 최종 방어(예외적으로 경계 잘림 실패 시)
                    if len(part.encode("utf-8")) > 32768:
                        part = part.encode("utf-8")[:32768].decode("utf-8", errors="ignore")

                    vec = hf_embed_text(tok, model, device, part, max_len=MAX_TOKENS)
                    
                    # 벡터 차원 검증
                    if len(vec) != emb_dim:
                        logger.error(f"[Ingest] 벡터 차원 불일치: 예상={emb_dim}, 실제={len(vec)}, 텍스트='{part[:50]}...'")
                        continue  # 이 벡터는 건너뛰기
                    
                    batch.append({
                        "embedding": vec.tolist(),
                        "path": str(rel_txt.as_posix()),
                        "chunk_idx": int(idx),
                        "task_type": task,
                        "security_level": lvl,
                        "doc_id": str(doc_id),
                        "version": int(version),
                        "page": int(page_num),  # 페이지 번호 추가
                        "text": part,
                    })
                    if len(batch) >= BATCH_SIZE:
                        client.insert(collection_name, batch)
                        total_inserted += len(batch)
                        batch = []

            # 2) 표 조각은 통합 파일에 이미 포함되어 있으므로 별도 인제스트하지 않음
            # (통합 파일을 파싱할 때 표도 함께 청크화되므로 중복 방지)

        if batch:
            client.insert(collection_name, batch)
            total_inserted += len(batch)

    # 인덱스/로딩 재보장 및 메타 저장(유추된 doc_id/version 반영)
    try:
        client.flush(collection_name)
    except Exception:
        pass
    _ensure_collection_and_index(client, emb_dim, metric="IP", collection_name=collection_name)

    # META에 doc_id/version 보정이 있었다면 저장
    try:
        META_JSON_PATH.write_text(json.dumps(meta, ensure_ascii=False, indent=2), encoding="utf-8")
    except Exception:
        pass

    return {
        "message": f"Ingest 완료(Milvus Server, collection={collection_name})",
        "inserted_chunks": int(total_inserted),
    }



# -------------------------------------------------
# 2-1) 단일 파일 인제스트(선택 작업유형)
# -------------------------------------------------
async def ingest_single_pdf(req: SinglePDFIngestRequest):
    try:
        from repository.documents import insert_workspace_document
    except Exception:
        insert_workspace_document = None

    file_path = Path(req.pdf_path)
    if not file_path.exists():
        return {"error": f"파일 경로를 찾을 수 없습니다: {file_path}"}

    if _ext(file_path) not in SUPPORTED_EXTS:
        return {"error": f"지원되지 않는 파일 형식입니다: {_ext(file_path)}"}

    # 메타 로드
    if META_JSON_PATH.exists():
        meta = json.loads(META_JSON_PATH.read_text(encoding="utf-8"))
    else:
        meta = {}

    # 추출
    text_all, table_blocks_all = _extract_any(file_path)

    # 보안 레벨 판정(본문+표)
    all_rules = get_security_level_rules_all()
    whole_for_level = text_all + "\n\n" + "\n\n".join(t.get("text","") for t in (table_blocks_all or []))
    sec_map = {task: _determine_level_for_task(whole_for_level, all_rules.get(task, {"maxLevel": 1, "levels": {}})) for task in TASK_TYPES}
    max_sec = max(sec_map.values()) if sec_map else 1
    sec_folder = f"securityLevel{int(max_sec)}"

    # 보관 및 텍스트 저장
    rel_file = Path(sec_folder) / file_path.name
    (LOCAL_DATA_ROOT / rel_file).parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(file_path, LOCAL_DATA_ROOT / rel_file)
    txt_path = EXTRACTED_TEXT_DIR / rel_file.with_suffix(".txt")
    txt_path.parent.mkdir(parents=True, exist_ok=True)
    txt_path.write_text(text_all, encoding="utf-8")

    doc_id, ver = _parse_doc_version(file_path.stem)
    meta[str(rel_file)] = {
        "chars": len(text_all),
        "lines": len(text_all.splitlines()),
        "preview": (_clean_text(text_all[:200].replace("\n", " ")) + "…") if text_all else "",
        "security_levels": sec_map,
        "doc_id": doc_id,
        "version": ver,
        "tables": table_blocks_all or [],
        "sourceExt": _ext(file_path),
    }
    META_JSON_PATH.parent.mkdir(parents=True, exist_ok=True)
    META_JSON_PATH.write_text(json.dumps(meta, ensure_ascii=False, indent=2), encoding="utf-8")

    # 인제스트
    settings = get_vector_settings()
    tok, model, device = await _get_or_load_embedder_async(settings["embeddingModel"])
    emb_dim = int(hf_embed_text(tok, model, device, "probe").shape[0])
    client = _client()
    _ensure_collection_and_index(client, emb_dim, metric="IP", collection_name=COLLECTION_NAME)

    s = get_vector_settings()
    MAX_TOKENS, OVERLAP = int(s["chunkSize"]), int(s["overlap"])

    # 기존 삭제
    try:
        client.delete(COLLECTION_NAME, filter=f"doc_id == '{doc_id}' && version <= {int(ver)}")
    except Exception:
        pass

    tasks = req.task_types or list(TASK_TYPES)
    chunks = chunk_text(text_all, max_tokens=MAX_TOKENS, overlap=OVERLAP)
    batch, cnt = [], 0

    for task in tasks:
        lvl = int(sec_map.get(task, 1))

        # 본문: VARCHAR 안전 분할
        for idx, c in enumerate(chunks):
            for part in _split_for_varchar_bytes(c):
                if len(part.encode("utf-8")) > 32768:
                    part = part.encode("utf-8")[:32768].decode("utf-8", errors="ignore")
                vec = hf_embed_text(tok, model, device, part, max_len=MAX_TOKENS)
                batch.append({
                    "embedding": vec.tolist(),
                    "path": str(rel_file.with_suffix(".txt")),
                    "chunk_idx": int(idx),
                    "task_type": task,
                    "security_level": lvl,
                    "doc_id": str(doc_id),
                    "version": int(ver),
                    "text": part,
                })
                if len(batch) >= 128:
                    client.insert(COLLECTION_NAME, batch)
                    cnt += len(batch)
                    batch = []

        # 표: VARCHAR 안전 분할
        base_idx = len(chunks)
        for t_i, t in enumerate(table_blocks_all or []):
            md = (t.get("text") or "").strip()
            if not md:
                continue
            page = int(t.get("page", 0))
            bbox = t.get("bbox") or []
            bbox_str = ",".join(str(x) for x in bbox) if bbox else ""
            table_text = f"[[TABLE page={page} bbox={bbox_str}]]\n{md}"

            for sub_j, part in enumerate(_split_for_varchar_bytes(table_text)):
                if len(part.encode("utf-8")) > 32768:
                    part = part.encode("utf-8")[:32768].decode("utf-8", errors="ignore")
                vec = hf_embed_text(tok, model, device, part, max_len=MAX_TOKENS)
                batch.append({
                    "embedding": vec.tolist(),
                    "path": str(rel_file.with_suffix(".txt")),
                    "chunk_idx": int(base_idx + t_i * 1000 + sub_j),
                    "task_type": task,
                    "security_level": lvl,
                    "doc_id": str(doc_id),
                    "version": int(ver),
                    "text": part,
                })
                if len(batch) >= 128:
                    client.insert(COLLECTION_NAME, batch)
                    cnt += len(batch)
                    batch = []

    if batch:
        client.insert(COLLECTION_NAME, batch)
        cnt += len(batch)

    try:
        client.flush(COLLECTION_NAME)
    except Exception:
        pass
    _ensure_collection_and_index(client, emb_dim, metric="IP", collection_name=COLLECTION_NAME)

    return {
        "message": f"단일 파일 인제스트 완료(Milvus Server) - {_ext(file_path)}",
        "doc_id": doc_id,
        "version": ver,
        "chunks": cnt,
        "sourceExt": _ext(file_path),
    }

async def ingest_specific_files_with_levels(
    uploads: Optional[List[Any]] = None,          # FastAPI UploadFile 리스트
    paths: Optional[List[str]] = None,            # 로컬 경로 리스트
    tasks: Optional[List[str]] = None,            # 없으면 모든 TASK_TYPES
    level_for_tasks: Optional[Dict[str, int]] = None,  # {"qna":2,"summary":1} 우선
    level: Optional[int] = None,                  # 공통 레벨. 위 map 있으면 무시
    collection_name: Optional[str] = None,
):
    if not uploads and not paths:
        return {"error": "대상 파일이 없습니다. uploads 또는 paths 중 하나는 필요합니다."}

    tasks_eff = [t for t in (tasks or TASK_TYPES) if t in TASK_TYPES]
    if not tasks_eff:
        return {"error": f"유효한 작업유형이 없습니다. 허용: {TASK_TYPES}"}

    lvl_map: Dict[str, int] = {}
    if level_for_tasks:
        for k, v in level_for_tasks.items():
            if k in TASK_TYPES:
                lvl_map[k] = max(1, int(v))
    elif level is not None:
        for t in tasks_eff:
            lvl_map[t] = max(1, int(level))

    # 업로드 저장(임시) + 경로 합치기
    run_id = uuid.uuid4().hex[:8]
    tmp_root = (VAL_SESSION_ROOT / "adhoc" / run_id).resolve()
    tmp_root.mkdir(parents=True, exist_ok=True)

    saved: List[Path] = []
    if uploads:
        for f in uploads:
            fname = Path(getattr(f, "filename", "uploaded")).name
            tmp_path = tmp_root / fname
            try:
                data = await f.read()
            except Exception:
                data = getattr(getattr(f, "file", None), "read", lambda: b"")()
            tmp_path.write_bytes(data or b"")
            saved.append(tmp_path)
    for p in (paths or []):
        pp = Path(str(p)).resolve()
        if pp.exists() and pp.is_file():
            saved.append(pp)

    if not saved:
        return {"error": "저장/유효성 검사 후 남은 파일이 없습니다."}

    # 임베더/컬렉션 준비
    settings = get_vector_settings()
    eff_model_key = settings["embeddingModel"]
    tok, model, device = await _get_or_load_embedder_async(eff_model_key)
    emb_dim = int(hf_embed_text(tok, model, device, "probe").shape[0])

    coll = collection_name or COLLECTION_NAME
    client = _client()
    _ensure_collection_and_index(client, emb_dim, metric="IP", collection_name=coll)

    MAX_TOKENS, OVERLAP = int(settings["chunkSize"]), int(settings["overlap"])

    processed, total = [], 0
    for src in saved:
        try:
            text, tables = _extract_any(src)

            # 레벨 결정(강제 > 규칙)
            if lvl_map:
                sec_map = {t: int(lvl_map.get(t, 1)) for t in tasks_eff}
            else:
                all_rules = get_security_level_rules_all()
                whole = text + "\n\n" + "\n\n".join(t.get("text", "") for t in (tables or []))
                sec_map = {
                    t: _determine_level_for_task(whole, all_rules.get(t, {"maxLevel": 1, "levels": {}}))
                    for t in tasks_eff
                }
            max_sec = max(sec_map.values()) if sec_map else 1

            # 스니펫 로딩용 텍스트 저장(메인과 분리: __adhoc__)
            rel_txt = Path("__adhoc__") / run_id / f"securityLevel{int(max_sec)}" / src.with_suffix(".txt").name
            abs_txt = EXTRACTED_TEXT_DIR / rel_txt
            abs_txt.parent.mkdir(parents=True, exist_ok=True)
            abs_txt.write_text(text, encoding="utf-8")

            # 문서 ID/버전
            doc_id, ver = _parse_doc_version(src.stem)

            # 기존 삭제
            try:
                client.delete(collection_name=coll, filter=f"doc_id == '{doc_id}' && version <= {int(ver)}")
            except Exception:
                pass

            # 본문
            chunks = chunk_text(text, max_tokens=MAX_TOKENS, overlap=OVERLAP)
            batch, cnt = [], 0
            for t in tasks_eff:
                lvl = int(sec_map.get(t, 1))

                for idx, c in enumerate(chunks):
                    for part in _split_for_varchar_bytes(c):
                        if len(part.encode("utf-8")) > 32768:
                            part = part.encode("utf-8")[:32768].decode("utf-8", errors="ignore")
                        vec = hf_embed_text(tok, model, device, part, max_len=MAX_TOKENS)
                        batch.append({
                            "embedding": vec.tolist(),
                            "path": str(rel_txt.as_posix()),
                            "chunk_idx": int(idx),
                            "task_type": t,
                            "security_level": lvl,
                            "doc_id": str(doc_id),
                            "version": int(ver),
                            "text": part,
                        })
                        if len(batch) >= 128:
                            client.insert(coll, batch); cnt += len(batch); batch = []

                # 표
                base_idx = len(chunks)
                for t_i, tb in enumerate(tables or []):
                    md = (tb.get("text") or "").strip()
                    if not md:
                        continue
                    page = int(tb.get("page", 0)); bbox = tb.get("bbox") or []
                    bbox_str = ",".join(str(x) for x in bbox) if bbox else ""
                    table_text = f"[[TABLE page={page} bbox={bbox_str}]]\n{md}"
                    for sub_j, part in enumerate(_split_for_varchar_bytes(table_text)):
                        if len(part.encode("utf-8")) > 32768:
                            part = part.encode("utf-8")[:32768].decode("utf-8", errors="ignore")
                        vec = hf_embed_text(tok, model, device, part, max_len=MAX_TOKENS)
                        batch.append({
                            "embedding": vec.tolist(),
                            "path": str(rel_txt.as_posix()),
                            "chunk_idx": int(base_idx + t_i * 1000 + sub_j),
                            "task_type": t,
                            "security_level": lvl,
                            "doc_id": str(doc_id),
                            "version": int(ver),
                            "text": part,
                        })
                        if len(batch) >= 128:
                            client.insert(coll, batch); cnt += len(batch); batch = []

            if batch:
                client.insert(coll, batch); cnt += len(batch)

            processed.append({
                "file": src.name, "doc_id": doc_id, "version": int(ver),
                "levels": sec_map, "chunks": cnt
            })
            total += cnt

        except Exception:
            logger.exception("[upload-and-ingest] failed: %s", src)

    try:
        client.flush(coll)
    except Exception:
        pass
    _ensure_collection_and_index(client, emb_dim, metric="IP", collection_name=coll)

    return {
        "message": "Upload & Ingest 완료",
        "collection": coll,
        "runId": run_id,
        "processed": processed,
        "inserted_chunks": int(total),
    }

# -------------------------------------------------
# 3) 검색 (vector / hybrid)
#   - task_type 필터 + security_level 제한
#   - hybrid: 벡터 topK*α 후보에 대해 간이 BM25 후처리 리랭크
# -------------------------------------------------
def _bm25_like_score(query: str, doc: str, k1: float = 1.2, b: float = 0.75) -> float:
    # 후보군 소규모 리랭크용 간단 BM25 대용(문서 집합이 작을 때만)
    # 토크나이징 매우 단순화(공백 기준)
    q_terms = [w for w in query.lower().split() if w]
    d_terms = [w for w in doc.lower().split() if w]
    if not q_terms or not d_terms:
        return 0.0
    d_len = len(d_terms)
    tf = Counter(d_terms)
    # IDF는 후보 집합 크기를 사용하기 어려워 고정치에 완화 가중
    score = 0.0
    avgdl = max(1.0, d_len)  # 후보 단일 문서 기준
    for t in set(q_terms):
        f = tf.get(t, 0)
        if f == 0:
            continue
        # 완화 IDF(상수): log(1 + 1/freq) 대신 상수 1.5 사용(경험적)
        idf = 1.5
        denom = f + k1 * (1 - b + b * (d_len / avgdl))
        score += idf * ((f * (k1 + 1)) / (denom if denom != 0 else 1))
    return float(score)


async def search_documents(req: RAGSearchRequest, search_type_override: Optional[str] = None,
                           collection_name: str = COLLECTION_NAME, rerank_top_n: Optional[int] = None) -> Dict:
    t0 = time.perf_counter()
    print(f"🔍 [Search] 검색 시작: query='{req.query}', topK={req.top_k}, rerank_topN={rerank_top_n}, task={req.task_type}")
    
    if req.task_type not in TASK_TYPES:
        return {
            "error": f"invalid task_type: {req.task_type}. choose one of {TASK_TYPES}"
        }

    settings = get_vector_settings()
    model_key = req.model or settings["embeddingModel"]
    raw_st = (search_type_override or settings.get("searchType") or "").lower()
    # alias normalization: 'semantic'/'sementic' -> 'vector'; default 'hybrid' if empty
    search_type = (raw_st.replace("semantic", "vector").replace("sementic", "vector") or "hybrid")

    tok, model, device = await _get_or_load_embedder_async(model_key)
    q_emb = hf_embed_text(tok, model, device, req.query)
    client = _client()
    _ensure_collection_and_index(client, emb_dim=len(q_emb), metric="IP")

    if COLLECTION_NAME not in client.list_collections():
        return {"error": "컬렉션이 없습니다. 먼저 데이터 저장(인제스트)을을 수행하세요."}

    # 공통 파라미터
    embedding_candidates = int(req.top_k)  # 임베딩에서 찾을 후보 개수
    final_results = int(rerank_top_n) if rerank_top_n is not None else 5  # 최종 반환 개수
    candidate = max(embedding_candidates, final_results * 2)  # 충분한 후보 확보
    filter_expr = f"task_type == '{req.task_type}' && security_level <= {int(req.user_level)}"

    def _dense_search(limit=candidate):
        return client.search(
            collection_name=COLLECTION_NAME,
            data=[q_emb.tolist()],
            anns_field="embedding",
            limit=int(limit),
            search_params={"metric_type": "IP", "params": {}},
            output_fields=["path", "chunk_idx", "task_type", "security_level", "doc_id", "text", "page"],
            filter=filter_expr,
        )

    def _sparse_search(query_text: str, limit=candidate):
        """
        BM25 Function(text->text_sparse)이 붙어 있으면 서버가 스파스 점수를 계산한다.
        최신 pymilvus에서는 anns_field='text_sparse', data=['쿼리 문자열'] 형태를 지원.
        일부 버전에서 미지원이면 예외 발생 → 폴백으로 빈 결과 반환.
        """
        try:
            return client.search(
                collection_name=COLLECTION_NAME,
                data=[query_text],
                anns_field="text_sparse",
                limit=int(limit),
                search_params={"params": {}},
                output_fields=["path", "chunk_idx", "task_type", "security_level", "doc_id", "text", "page"],
                filter=filter_expr,
            )
        except Exception as e:
            logger.warning(f"[Milvus] sparse search unavailable: {e}")
            return [[]]

    def _load_snippet(
        path: str, cidx: int, max_tokens: int = 512, overlap: int = 64
    ) -> str:
        file_path = EXTRACTED_TEXT_DIR / path

        logger.debug(f"\n###########################\nfile_path: {file_path}")
        try:
            full_txt = file_path.read_text(encoding="utf-8")
        except Exception as exc:
            logger.warning(f"[Milvus] snippet 파일 로드 실패: {file_path} ({exc})")
            full_txt = ""
        if not full_txt:
            # 필요하면 ent_text로부터 넘어온 값이나 최소 안내 문구 반환
            return ""

        words = full_txt.split()
        if not words:
            return ""

        window = max_tokens - overlap
        if window <= 0:
            window = max_tokens

        start = max(0, cidx * window)
        snippet = " ".join(words[start:start + max_tokens]).strip()
        if snippet:
            return snippet

        # fallback: 청크 범위가 벗어나면 처음 구간이라도 리턴
        return " ".join(words[:max_tokens]).strip()

    # === 분기: 검색 방식 ===
    hits_raw = []
    TABLE_MARK = "[[TABLE"
    
    if search_type == "vector":
        results = _dense_search(limit=candidate)
        for hit in results[0]:
            if isinstance(hit, dict):
                ent = hit.get("entity", {})
                ent_text = ent.get("text")  # ★ 추가
                path = ent.get("path")
                cidx = int(ent.get("chunk_idx", 0))
                ttype = ent.get("task_type")
                lvl = int(ent.get("security_level", 1))
                doc_id = ent.get("doc_id")
                page = int(ent.get("page", 0))  # 페이지 정보 추출
                score_vec = float(hit.get("distance", 0.0))
            else:
                ent = hit.entity
                ent_text = getattr(ent, "get", lambda _k: None)("text") if hasattr(ent, "get") else None
                path = hit.entity.get("path")
                cidx = int(hit.entity.get("chunk_idx", 0))
                ttype = hit.entity.get("task_type")
                lvl = int(hit.entity.get("security_level", 1))
                doc_id = hit.entity.get("doc_id")
                page = int(hit.entity.get("page", 0))  # 페이지 정보 추출
                score_vec = float(hit.score)
            
            # 스니펫 결정 로직: 표면 저장된 텍스트 그대로, 아니면 기존 로직
            if isinstance(ent_text, str) and ent_text.startswith(TABLE_MARK):
                snippet = ent_text  # ★ 표는 저장된 마크다운 그대로
            else:
                snippet = _load_snippet(path, cidx)
                
            hits_raw.append({
                "path": path, "chunk_idx": cidx, "task_type": ttype,
                "security_level": lvl, "doc_id": doc_id, "page": page,
                "score_vec": score_vec, "score_sparse": 0.0, "snippet": snippet
            })
    else:
        # hybrid / bm25: 덴스 + 스파스 각각 검색
        res_dense = _dense_search(limit=candidate)
        res_sparse = _sparse_search(req.query, limit=candidate)

        def _collect(res, is_dense: bool):
            out = []
            for hit in (res[0] if res and len(res) > 0 else []):
                if isinstance(hit, dict):
                    ent = hit.get("entity", {})
                    ent_text = ent.get("text")  # ★ 추가
                    path = ent.get("path")
                    cidx = int(ent.get("chunk_idx", 0))
                    ttype = ent.get("task_type")
                    lvl = int(ent.get("security_level", 1))
                    doc_id = ent.get("doc_id")
                    page = int(ent.get("page", 0))  # 페이지 정보 추출
                    score = float(hit.get("distance", 0.0))
                else:
                    ent = hit.entity
                    ent_text = getattr(ent, "get", lambda _k: None)("text") if hasattr(ent, "get") else None
                    path = hit.entity.get("path")
                    cidx = int(hit.entity.get("chunk_idx", 0))
                    ttype = hit.entity.get("task_type")
                    lvl = int(hit.entity.get("security_level", 1))
                    doc_id = hit.entity.get("doc_id")
                    page = int(hit.entity.get("page", 0))  # 페이지 정보 추출
                    score = float(hit.score)
                out.append(((path, cidx, ttype, lvl, doc_id, page, ent_text), score))  # ★ page, ent_text 추가
            return out

        dense_list = _collect(res_dense, True)
        sparse_list = _collect(res_sparse, False)

        # RRF 결합 폴백
        # key_short는 (path, cidx, ttype, lvl, doc_id)로 설정하여 같은 청크는 하나만 나타나도록 함
        # 페이지 정보는 별도로 저장하여 페이로드에 포함
        rrf: dict[tuple, float] = {}
        text_map: dict[tuple, str] = {}  # ★ 텍스트 매핑 추가
        page_map: dict[tuple, int] = {}  # ★ 페이지 매핑 추가
        score_map: dict[tuple, tuple[float, float]] = {}  # (score_vec, score_sparse) 저장
        K = 60.0
        for rank, (key, score) in enumerate(dense_list, start=1):
            key_short = key[:5]  # (path, cidx, ttype, lvl, doc_id) - page 제외
            rrf[key_short] = rrf.get(key_short, 0.0) + 1.0 / (K + rank)
            if len(key) > 6:  # ent_text가 있으면 저장
                text_map[key_short] = key[6]
            if len(key) > 5:  # page가 있으면 저장 (첫 번째로 발견된 페이지 사용)
                if key_short not in page_map:
                    page_map[key_short] = key[5]
            # 점수 저장 (dense) - 최대값 유지
            if key_short not in score_map:
                score_map[key_short] = (score, 0.0)
            else:
                score_map[key_short] = (max(score_map[key_short][0], score), score_map[key_short][1])
        
        for rank, (key, score) in enumerate(sparse_list, start=1):
            key_short = key[:5]  # (path, cidx, ttype, lvl, doc_id) - page 제외
            rrf[key_short] = rrf.get(key_short, 0.0) + 1.0 / (K + rank)
            if len(key) > 6:  # ent_text가 있으면 저장
                text_map[key_short] = key[6]
            if len(key) > 5:  # page가 있으면 저장 (아직 없으면 저장)
                if key_short not in page_map:
                    page_map[key_short] = key[5]
            # 점수 저장 (sparse) - 최대값 유지
            if key_short not in score_map:
                score_map[key_short] = (0.0, score)
            else:
                score_map[key_short] = (score_map[key_short][0], max(score_map[key_short][1], score))

        merged = sorted(rrf.items(), key=lambda x: x[1], reverse=True)[:candidate]
        for (path, cidx, ttype, lvl, doc_id), fused in merged:
            # 스니펫 결정 로직
            ent_text = text_map.get((path, cidx, ttype, lvl, doc_id))
            if isinstance(ent_text, str) and ent_text.startswith(TABLE_MARK):
                snippet = ent_text  # ★ 표는 저장된 마크다운 그대로
            else:
                snippet = _load_snippet(path, cidx)
            
            page_num = page_map.get((path, cidx, ttype, lvl, doc_id), 0)
            s_vec, s_spa = score_map.get((path, cidx, ttype, lvl, doc_id), (0.0, 0.0))
            hits_raw.append({
                "path": path, "chunk_idx": cidx, "task_type": ttype,
                "security_level": lvl, "doc_id": doc_id, "page": int(page_num),
                "score_vec": float(s_vec), "score_sparse": float(s_spa),
                "score_fused": float(fused),
                "snippet": snippet
            })

    # 검색 결과 상태 로그
    logger.info(f"📊 [Search] 벡터/BM25 검색 완료: 후보 {len(hits_raw)}개 발견")
    if hits_raw:
        logger.info(f"📊 [Search] 첫 번째 후보: doc_id={hits_raw[0].get('doc_id')}, path={hits_raw[0].get('path')}")

    rerank_candidates = []
    for hit in hits_raw:
        snippet = hit.get("snippet", "")
        if not snippet:
            continue
        rerank_candidates.append(
            {
                "text": snippet,
                "score": float(hit.get("score_fused", hit.get("score_vec", hit.get("score_sparse", 0.0)) or 0.0)),
                "doc_id": hit.get("doc_id"),
                "title": hit.get("doc_id") or hit.get("path") or "snippet",
                "source": "milvus",
                "metadata": hit,
            }
        )

    if rerank_candidates:
        reranked = rerank_snippets(rerank_candidates, query=req.query, top_n=final_results)
        hits_sorted = []
        for res in reranked:
            original = res.metadata or {}
            hits_sorted.append(
                {
                    "score": float(res.score),
                    "path": original.get("path"),
                    "chunk_idx": int(original.get("chunk_idx", 0)),
                    "task_type": original.get("task_type"),
                    "security_level": int(original.get("security_level", 1)),
                    "doc_id": original.get("doc_id"),
                    "page": int(original.get("page", 0)),
                    "snippet": res.text,
                }
            )
    else:
        hits_sorted = sorted(
            hits_raw,
            key=lambda x: x.get("score_fused", x.get("score_vec", x.get("score_sparse", 0.0))),
            reverse=True,
        )[:final_results]

    # 리랭크 결과 로그 출력
    if hits_sorted:
        top_hit = hits_sorted[0]
        logger.info(f"✨ [Rerank] 완료! 최고 점수: {top_hit.get('score', 0):.4f}")
        logger.info(f"🏆 [Rerank] 최고 스니펫 (doc_id: {top_hit.get('doc_id', 'unknown')}): {top_hit.get('snippet', '')[:100]}...")

    # 프롬프트 컨텍스트 생성
    context = "\n---\n".join(h["snippet"] for h in hits_sorted if h.get("snippet"))
    prompt = f"사용자 질의: {req.query}\n:\n{context}\n\n위 내용을 바탕으로 응답을 생성해 주세요."

    elapsed = round(time.perf_counter() - t0, 4)

    # query_logs 삭제: INSERT 제거
    return {
        "elapsed_sec": elapsed,
        "settings_used": {"model": model_key, "searchType": search_type},
        "hits": [
            {
                "score": float(h["score"]),
                "path": h["path"],
                "chunk_idx": int(h["chunk_idx"]),
                "task_type": h["task_type"],
                "security_level": int(h["security_level"]),
                "doc_id": h.get("doc_id"),
                "page": int(h.get("page", 0)),  # 페이지 정보 추가
                "snippet": h["snippet"],
            }
            for h in hits_sorted
        ],
        "prompt": prompt,
    }




async def execute_search(
    question: str,
    top_k: int = 20,   # 임베딩 후보 개수
    rerank_top_n: int = 5,    # 최종 반환 개수  
    security_level: int = 1,
    source_filter: Optional[List[str]] = None,
    task_type: str = "qna",
    model_key: Optional[str] = None,
    search_type: Optional[str] = None,
) -> Dict:
    print(f"⭐ [ExecuteSearch] 함수 호출: question='{question}', topK={top_k}, rerank_topN={rerank_top_n}")
    req = RAGSearchRequest(
        query=question,
        top_k=top_k,
        user_level=security_level,
        task_type=task_type,
        model=model_key,
    )
    logger.info(f"📞 [ExecuteSearch] search_documents 호출 전: req 생성 완료")
    res = await search_documents(req, search_type_override=search_type, rerank_top_n=rerank_top_n)
    logger.info(f"📞 [ExecuteSearch] search_documents 호출 후: 결과 hits 수={len(res.get('hits', []))}")
    # Build check_file BEFORE optional source_filter so it reflects original candidates
    check_files: List[str] = []
    logger.debug(f"\n###########################\nres: {res}")
    try:
        for h in res.get("hits", []):
            # Prefer doc_id when available; fallback to path-derived filename
            doc_id_val = h.get("doc_id")
            if doc_id_val:
                check_files.append(f"{str(doc_id_val)}.pdf")
                continue
            p = Path(h.get("path", ""))
            if str(p):
                check_files.append(p.with_suffix(".pdf").name)
    except Exception:
        pass

    if source_filter and "hits" in res:
        names = {Path(n).stem for n in source_filter}
        res["hits"] = [h for h in res["hits"] if Path(h["path"]).stem in names]

    res["check_file"] = sorted(list(set(check_files)))
    return res

# -------------------------------------------------
# 4) 관리 유틸
# -------------------------------------------------
async def delete_db():
    # 모델 캐시 클리어
    _invalidate_embedder_cache()

    client = _client()
    cols = client.list_collections()
    for c in cols:
        client.drop_collection(c)
    return {"message": "삭제 완료(Milvus Server)", "dropped_collections": cols}

async def list_indexed_files(
    limit: int = 16384,
    offset: int = 0,
    query: Optional[str] = None,
    task_type: Optional[str] = None,
):
    limit = max(1, min(limit, 16384))
    client = _client()
    if COLLECTION_NAME not in client.list_collections():
        return []

    # 메타 로드(원본 확장자 복원용)
    try:
        meta = json.loads(META_JSON_PATH.read_text(encoding="utf-8"))
    except Exception:
        meta = {}

    flt = ""
    if task_type and task_type in TASK_TYPES:
        flt = f"task_type == '{task_type}'"
    try:
        rows = client.query(
            collection_name=COLLECTION_NAME,
            filter=flt,
            output_fields=["path", "chunk_idx", "security_level", "task_type"],
            limit=limit,
            offset=offset,
            consistency_level="Strong",
        )
    except Exception:
        rows = []

    counts: Dict[Tuple[str, str], int] = defaultdict(int)
    level_map: Dict[Tuple[str, str], int] = {}
    for r in rows:
        path = r.get("path") if isinstance(r, dict) else r["path"]
        ttype = r.get("task_type") if isinstance(r, dict) else r["task_type"]
        lvl = int((r.get("security_level") if isinstance(r, dict) else r["security_level"]) or 1)
        key = (path, ttype)
        counts[key] += 1
        level_map.setdefault(key, lvl)

    items = []
    for (path, ttype), cnt in counts.items():
        txt_rel = Path(path)

        # 메타에서 원래 확장자를 복원
        cands = [txt_rel.with_suffix(ext).as_posix() for ext in SUPPORTED_EXTS]
        meta_key = next((k for k in cands if k in meta), None)
        if meta_key:
            source_ext = meta.get(meta_key, {}).get("sourceExt") or Path(meta_key).suffix
            orig_rel = txt_rel.with_suffix(source_ext)
        else:
            # 폴백(구버전 데이터): pdf 가정
            orig_rel = txt_rel.with_suffix(".pdf")

        file_name = orig_rel.name
        file_path = str(orig_rel)

        txt_abs = EXTRACTED_TEXT_DIR / txt_rel
        try:
            stat = txt_abs.stat()
            size = stat.st_size
            indexed_at = now_kst_string()
        except FileNotFoundError:
            size = None
            indexed_at = None
        items.append(
            {
                "taskType": ttype,
                "fileName": file_name,
                "filePath": file_path,
                "chunkCount": int(cnt),
                "indexedAt": indexed_at,
                "fileSize": size,
                "securityLevel": int(level_map.get((path, ttype), 1)),
            }
        )

    if query:
        q = str(query)
        items = [it for it in items if q in it["fileName"]]
    return items

async def delete_files_by_names(file_names: List[str], task_type: Optional[str] = None, collection_name: str = COLLECTION_NAME):
    """
    파일명(= doc_id stem) 배열을 받아 벡터 DB에서 삭제.
    - task_type 가 None 이면 모든 작업유형(doc_gen/summary/qna)에서 삭제 (기존 동작과 동일)
    - task_type 가 지정되면 해당 작업유형 레코드만 삭제
    """
    if not file_names:
        return {"deleted": 0, "requested": 0}

    try:
        from repository.documents import delete_workspace_documents_by_filenames
    except Exception:
        delete_workspace_documents_by_filenames = None

    client = _client()
    if COLLECTION_NAME not in client.list_collections():
        deleted_sql = None
        if delete_workspace_documents_by_filenames:
            deleted_sql = delete_workspace_documents_by_filenames(file_names)
        return {"deleted": 0, "deleted_sql": deleted_sql, "requested": len(file_names)}

    # 로드 보장
    try:
        client.load_collection(collection_name=COLLECTION_NAME)
    except Exception:
        pass

    # 유효한 task_type 인지 검증
    task_filter = ""
    if task_type:
        if task_type not in TASK_TYPES:
            return {
                "deleted": 0,
                "requested": len(file_names),
                "error": f"invalid taskType: {task_type}",
            }
        task_filter = f" && task_type == '{task_type}'"

    deleted_total = 0
    per_file: dict[str, int] = {}

    for name in file_names:
        stem = Path(name).stem
        # Align fileName -> doc_id by stripping version suffix if present
        try:
            base_id, _ver = _parse_doc_version(stem)
        except Exception:
            base_id = stem
        try:
            # doc_id == 'stem' [&& task_type == 'xxx']
            filt = f"doc_id == '{base_id}'{task_filter}"
            client.delete(collection_name=COLLECTION_NAME, filter=filt)
            deleted_total += 1
            per_file[name] = per_file.get(name, 0) + 1
        except Exception:
            logger.exception("Failed to delete from Milvus for file: %s", name)
            per_file[name] = per_file.get(name, 0)

    # Ensure deletion is visible to subsequent queries (file lists/overview)
    try:
        client.flush(COLLECTION_NAME)
    except Exception:
        logger.exception("Failed to flush Milvus after deletion")
    # Force reload to avoid any stale cache/state on the server side
    try:
        client.release_collection(collection_name=COLLECTION_NAME)
    except Exception:
        pass
    try:
        client.load_collection(collection_name=COLLECTION_NAME)
    except Exception:
        logger.exception("Failed to reload collection after deletion")

    deleted_sql = None
    if delete_workspace_documents_by_filenames:
        try:
            # SQL은 작업유형 구분이 없다고 가정(기존 그대로)
            deleted_sql = delete_workspace_documents_by_filenames(file_names)
        except Exception:
            logger.exception("Failed to delete workspace documents in SQL")
            deleted_sql = None

    return {
        "deleted": deleted_total,  # 요청 파일 기준 성공 건수(작업유형 기준 단순 카운트)
        "deleted_sql": deleted_sql,
        "requested": len(file_names),
        "taskType": task_type,
        "perFile": per_file,  # 파일별 처리현황
    }


async def list_indexed_files_overview(collection_name: str = COLLECTION_NAME):
    items = await list_indexed_files(limit=16384, offset=0, query=None, task_type=None)
    # agg: task_type -> level -> count
    agg: Dict[str, Dict[int, int]] = defaultdict(lambda: defaultdict(int))
    for it in items:
        agg[it["taskType"]][int(it["securityLevel"])] += it["chunkCount"]
    # 보기 좋게 변환
    overview = {
        t: {str(lv): agg[t][lv] for lv in sorted(agg[t].keys())} for t in agg.keys()
    }
    return {"overview": overview, "items": items}


# --- add: 세션 인덱스 로드/저장 + 컬렉션명 ---
def _load_sessions_index() -> dict:
    VAL_SESSION_ROOT.mkdir(parents=True, exist_ok=True)
    if SESSIONS_INDEX_PATH.exists():
        try:
            return json.loads(SESSIONS_INDEX_PATH.read_text(encoding="utf-8"))
        except Exception:
            logger.exception("failed to read sessions index")
    return {}

def _save_sessions_index(idx: dict):
    VAL_SESSION_ROOT.mkdir(parents=True, exist_ok=True)
    SESSIONS_INDEX_PATH.write_text(json.dumps(idx, ensure_ascii=False, indent=2), encoding="utf-8")

def _session_collection_name(sid: str) -> str:
    return f"{COLLECTION_NAME}__sess__{sid}"

# --- add: 세션 생성/조회/삭제 ---
def create_test_session() -> Dict:
    idx = _load_sessions_index()
    sid = uuid.uuid4().hex[:12]
    sess_dir = (VAL_SESSION_ROOT / sid).resolve()
    sess_dir.mkdir(parents=True, exist_ok=True)

    meta = {
        "sid": sid,
        "dir": str(sess_dir),
        "collection": _session_collection_name(sid),
        "createdAt": now_kst_string(),
    }
    idx[sid] = meta
    _save_sessions_index(idx)
    return meta

def get_test_session(sid: str) -> Optional[Dict]:
    idx = _load_sessions_index()
    return idx.get(sid)

async def drop_test_session(sid: str) -> Dict:
    idx = _load_sessions_index()
    meta = idx.pop(sid, None)
    if not meta:
        return {"success": False, "error": "invalid sid"}

    # 1) Milvus 컬렉션 드롭
    try:
        client = _client()
        coll = meta.get("collection") or _session_collection_name(sid)
        if coll in client.list_collections():
            client.drop_collection(coll)
    except Exception:
        logger.exception("[test-session] drop collection failed")

    # 2) 로컬 디렉토리 삭제(업로드 PDF)
    try:
        p = Path(meta["dir"])
        if p.exists():
            shutil.rmtree(p, ignore_errors=True)
    except Exception:
        logger.exception("[test-session] remove session dir failed")

    # 3) 세션 텍스트 폴더(EXTRACTED_TEXT_DIR/__sessions__/sid) 삭제
    try:
        sess_txt = EXTRACTED_TEXT_DIR / "__sessions__" / sid
        if sess_txt.exists():
            shutil.rmtree(sess_txt, ignore_errors=True)
    except Exception:
        logger.exception("[test-session] remove session texts failed")

    _save_sessions_index(idx)
    return {"success": True, "sid": sid, "dropped": True}

# --- add: ingest_test_pdfs ---
async def ingest_test_pdfs(sid: str, pdf_paths: List[str], task_types: Optional[List[str]] = None):
    meta = get_test_session(sid)
    if not meta:
        return {"error": "invalid sid"}

    settings = get_vector_settings()
    eff_model_key = settings["embeddingModel"]
    tok, model, device = await _get_or_load_embedder_async(eff_model_key)
    emb_dim = int(hf_embed_text(tok, model, device, "probe").shape[0])

    client = _client()
    coll = meta.get("collection") or _session_collection_name(sid)
    _ensure_collection_and_index_for(client, collection_name=coll, emb_dim=emb_dim, metric="IP")

    MAX_TOKENS, OVERLAP = int(settings["chunkSize"]), int(settings["overlap"])

    all_rules = get_security_level_rules_all()
    sess_txt_root = EXTRACTED_TEXT_DIR / "__sessions__" / sid
    sess_txt_root.mkdir(parents=True, exist_ok=True)

    tasks = task_types or list(TASK_TYPES)
    total = 0

    for p in pdf_paths:
        file_path = Path(str(p))
        if _ext(file_path) not in SUPPORTED_EXTS:
            logger.warning(f"[test-ingest] Unsupported file type: {file_path}")
            continue

        try:
            file_text, table_blocks_all = _extract_any(file_path)
        except Exception:
            logger.exception("[test-ingest] read failed: %s", p)
            continue

        whole_for_level = file_text + "\n\n" + "\n\n".join(t.get("text","") for t in (table_blocks_all or []))
        sec_map = {t: _determine_level_for_task(whole_for_level, all_rules.get(t, {"maxLevel": 1, "levels": {}})) for t in TASK_TYPES}
        max_sec = max(sec_map.values()) if sec_map else 1
        sec_folder = f"securityLevel{int(max_sec)}"

        rel_txt = Path("__sessions__") / sid / sec_folder / file_path.with_suffix(".txt").name
        abs_txt = EXTRACTED_TEXT_DIR / rel_txt
        abs_txt.parent.mkdir(parents=True, exist_ok=True)
        abs_txt.write_text(file_text, encoding="utf-8")

        stem = file_path.stem
        doc_id, ver = _parse_doc_version(stem)

        try:
            client.delete(coll, filter=f"doc_id == '{doc_id}' && version <= {int(ver)}")
        except Exception:
            pass

        chunks = chunk_text(file_text, max_tokens=MAX_TOKENS, overlap=OVERLAP)
        batch: List[Dict] = []

        for t in tasks:
            lvl = int(sec_map.get(t, 1))

            # 본문: VARCHAR 안전 분할
            for idx, c in enumerate(chunks):
                for part in _split_for_varchar_bytes(c):
                    if len(part.encode("utf-8")) > 32768:
                        part = part.encode("utf-8")[:32768].decode("utf-8", errors="ignore")
                    vec = hf_embed_text(tok, model, device, part, max_len=MAX_TOKENS)
                    batch.append({
                        "embedding": vec.tolist(),
                        "path": str(rel_txt.as_posix()),
                        "chunk_idx": int(idx),
                        "task_type": t,
                        "security_level": lvl,
                        "doc_id": str(doc_id),
                        "version": int(ver),
                        "text": part,
                    })
                    if len(batch) >= 128:
                        client.insert(coll, batch)
                        total += len(batch)
                        batch = []

            # 표: VARCHAR 안전 분할
            base_idx = len(chunks)
            for t_i, table in enumerate(table_blocks_all or []):
                md = (table.get("text") or "").strip()
                if not md:
                    continue
                page = int(table.get("page", 0))
                bbox = table.get("bbox") or []
                bbox_str = ",".join(str(x) for x in bbox) if bbox else ""
                table_text = f"[[TABLE page={page} bbox={bbox_str}]]\n{md}"

                for sub_j, part in enumerate(_split_for_varchar_bytes(table_text)):
                    if len(part.encode("utf-8")) > 32768:
                        part = part.encode("utf-8")[:32768].decode("utf-8", errors="ignore")
                    vec = hf_embed_text(tok, model, device, part, max_len=MAX_TOKENS)
                    batch.append({
                        "embedding": vec.tolist(),
                        "path": str(rel_txt.as_posix()),
                        "chunk_idx": int(base_idx + t_i * 1000 + sub_j),
                        "task_type": t,
                        "security_level": lvl,
                        "doc_id": str(doc_id),
                        "version": int(ver),
                        "text": part,
                    })
                    if len(batch) >= 128:
                        client.insert(coll, batch)
                        total += len(batch)
                        batch = []

        if batch:
            client.insert(coll, batch)
            total += len(batch)

    try:
        client.flush(coll)
    except Exception:
        pass
    _ensure_collection_and_index_for(client, collection_name=coll, emb_dim=emb_dim, metric="IP")

    return {"message": "세션 인제스트 완료", "sid": sid, "inserted_chunks": total}


# --- add: search_documents_test ---
async def search_documents_test(req: RAGSearchRequest, sid: str, search_type_override: Optional[str] = None, rerank_top_n: Optional[int] = None) -> Dict:
    """
    세션 전용 컬렉션에서만 검색 (기존 search_documents의 세션 버전)
    """
    meta = get_test_session(sid)
    if not meta:
        return {"error": "invalid sid"}

    t0 = time.perf_counter()
    if req.task_type not in TASK_TYPES:
        return {"error": f"invalid task_type: {req.task_type}. choose one of {TASK_TYPES}"}

    settings = get_vector_settings()
    model_key = req.model or settings["embeddingModel"]
    raw_st = (search_type_override or settings.get("searchType") or "").lower()
    search_type = (raw_st.replace("semantic", "vector").replace("sementic", "vector") or "hybrid")

    tok, model, device = await _get_or_load_embedder_async(model_key)
    q_emb = hf_embed_text(tok, model, device, req.query)

    client = _client()
    coll = meta.get("collection") or _session_collection_name(sid)
    _ensure_collection_and_index_for(client, collection_name=coll, emb_dim=len(q_emb), metric="IP")
    if coll not in client.list_collections():
        return {"error": "세션 컬렉션이 없습니다. 먼저 인제스트 하세요."}

    embedding_candidates = int(req.top_k)  # 임베딩에서 찾을 후보 개수
    final_results = int(rerank_top_n) if rerank_top_n is not None else 5  # 최종 반환 개수
    candidate = max(embedding_candidates, final_results * 2)  # 충분한 후보 확보
    filter_expr = f"task_type == '{req.task_type}' && security_level <= {int(req.user_level)}"

    def _dense_search(limit=candidate):
        return client.search(
            collection_name=coll,
            data=[q_emb.tolist()],
            anns_field="embedding",
            limit=int(limit),
            search_params={"metric_type": "IP", "params": {}},
            output_fields=["path", "chunk_idx", "task_type", "security_level", "doc_id", "text"],
            filter=filter_expr,
        )

    def _sparse_search(query_text: str, limit=candidate):
        try:
            return client.search(
                collection_name=coll,
                data=[query_text],
                anns_field="text_sparse",
                limit=int(limit),
                search_params={"params": {}},
                output_fields=["path", "chunk_idx", "task_type", "security_level", "doc_id", "text"],
                filter=filter_expr,
            )
        except Exception as e:
            logger.warning(f"[Milvus] sparse search unavailable(test): {e}")
            return [[]]

    def _load_snippet_for_session(path: str, cidx: int, max_tokens: int = settings["chunkSize"], overlap: int = settings["overlap"]) -> str:
        # path는 EXTRACTED_TEXT_DIR 기준 상대경로("__sessions__/sid/...")로 저장되어 있음
        try:
            full_txt = (EXTRACTED_TEXT_DIR / path).read_text(encoding="utf-8")
        except Exception:
            return ""
        words = full_txt.split()
        if not words:
            return ""
        start = cidx * (max_tokens - overlap)
        snippet = " ".join(words[start : start + max_tokens]).strip()
        return snippet or " ".join(words[:max_tokens]).strip()

    hits_raw = []
    TABLE_MARK = "[[TABLE"
    
    if search_type == "vector":
        res = _dense_search(limit=candidate)
        for hit in res[0]:
            if isinstance(hit, dict):
                ent = hit.get("entity", {})
                ent_text = ent.get("text")  # ★ 추가
                path = ent.get("path")
                cidx = int(ent.get("chunk_idx", 0))
                ttype = ent.get("task_type")
                lvl = int(ent.get("security_level", 1))
                doc_id = ent.get("doc_id")
                score_vec = float(hit.get("distance", 0.0))
            else:
                ent = hit.entity
                ent_text = getattr(ent, "get", lambda _k: None)("text") if hasattr(ent, "get") else None
                path = hit.entity.get("path")
                cidx = int(hit.entity.get("chunk_idx", 0))
                ttype = hit.entity.get("task_type")
                lvl = int(hit.entity.get("security_level", 1))
                doc_id = hit.entity.get("doc_id")
                score_vec = float(hit.score)
            
            # 스니펫 결정 로직: 표면 저장된 텍스트 그대로, 아니면 기존 로직
            if isinstance(ent_text, str) and ent_text.startswith(TABLE_MARK):
                snippet = ent_text  # ★ 표는 저장된 마크다운 그대로
            else:
                snippet = _load_snippet_for_session(path, cidx)
                
            hits_raw.append({
                "path": path, "chunk_idx": cidx, "task_type": ttype, "security_level": lvl,
                "doc_id": doc_id, "score_vec": score_vec, "score_sparse": 0.0, "snippet": snippet
            })
    else:
        res_dense = _dense_search(limit=candidate)
        res_sparse = _sparse_search(req.query, limit=candidate)

        def _collect(res):
            out = []
            for hit in (res[0] if res and len(res) > 0 else []):
                if isinstance(hit, dict):
                    ent = hit.get("entity", {})
                    ent_text = ent.get("text")  # ★ 추가
                    path = ent.get("path")
                    cidx = int(ent.get("chunk_idx", 0))
                    ttype = ent.get("task_type")
                    lvl = int(ent.get("security_level", 1))
                    doc_id = ent.get("doc_id")
                    score = float(hit.get("distance", 0.0))
                else:
                    ent = hit.entity
                    ent_text = getattr(ent, "get", lambda _k: None)("text") if hasattr(ent, "get") else None
                    path = hit.entity.get("path")
                    cidx = int(hit.entity.get("chunk_idx", 0))
                    ttype = hit.entity.get("task_type")
                    lvl = int(hit.entity.get("security_level", 1))
                    doc_id = hit.entity.get("doc_id")
                    score = float(hit.score)
                out.append(((path, cidx, ttype, lvl, doc_id, ent_text), score))  # ★ ent_text 추가
            return out

        dense_list = _collect(res_dense)
        sparse_list = _collect(res_sparse)
        rrf: Dict[tuple, float] = {}
        text_map: Dict[tuple, str] = {}  # ★ 텍스트 매핑 추가
        K = 60.0
        for rank, (key, _s) in enumerate(dense_list, start=1):
            key_short = key[:5]  # (path, cidx, ttype, lvl, doc_id)
            rrf[key_short] = rrf.get(key_short, 0.0) + 1.0 / (K + rank)
            if len(key) > 5:  # ent_text가 있으면 저장
                text_map[key_short] = key[5]
        for rank, (key, _s) in enumerate(sparse_list, start=1):
            key_short = key[:5]  # (path, cidx, ttype, lvl, doc_id)
            rrf[key_short] = rrf.get(key_short, 0.0) + 1.0 / (K + rank)
            if len(key) > 5:  # ent_text가 있으면 저장
                text_map[key_short] = key[5]
        merged = sorted(rrf.items(), key=lambda x: x[1], reverse=True)[:candidate]
        for (path, cidx, ttype, lvl, doc_id), fused in merged:
            # 스니펫 결정 로직
            ent_text = text_map.get((path, cidx, ttype, lvl, doc_id))
            if isinstance(ent_text, str) and ent_text.startswith(TABLE_MARK):
                snippet = ent_text  # ★ 표는 저장된 마크다운 그대로
            else:
                snippet = _load_snippet_for_session(path, cidx)
            
            s_vec = next((s for (k, s) in dense_list if k[:5] == (path, cidx, ttype, lvl, doc_id)), 0.0)
            s_spa = next((s for (k, s) in sparse_list if k[:5] == (path, cidx, ttype, lvl, doc_id)), 0.0)
            hits_raw.append({
                "path": path, "chunk_idx": cidx, "task_type": ttype, "security_level": lvl, "doc_id": doc_id,
                "score_vec": float(s_vec), "score_sparse": float(s_spa), "score_fused": float(fused), "snippet": snippet
            })

    rerank_candidates = []
    for hit in hits_raw:
        snippet = hit.get("snippet", "")
        if not snippet:
            continue
        rerank_candidates.append(
            {
                "text": snippet,
                "score": float(hit.get("score_fused", hit.get("score_vec", hit.get("score_sparse", 0.0)) or 0.0)),
                "doc_id": hit.get("doc_id"),
                "title": hit.get("doc_id") or hit.get("path") or "snippet",
                "source": "milvus",
                "metadata": hit,
            }
        )

    if rerank_candidates:
        reranked = rerank_snippets(rerank_candidates, query=req.query, top_n=final_results)
        hits_sorted = []
        for res in reranked:
            original = res.metadata or {}
            hits_sorted.append(
                {
                    "score": float(res.score),
                    "path": original.get("path"),
                    "chunk_idx": int(original.get("chunk_idx", 0)),
                    "task_type": original.get("task_type"),
                    "security_level": int(original.get("security_level", 1)),
                    "doc_id": original.get("doc_id"),
                    "snippet": res.text,
                }
            )
    else:
        hits_sorted = sorted(
            hits_raw,
            key=lambda x: x.get("score_fused", x.get("score_vec", x.get("score_sparse", 0.0))),
            reverse=True,
        )[:final_results]

    # 리랭크 결과 로그 출력 (테스트 세션)
    if hits_sorted:
        top_hit = hits_sorted[0]
        logger.info(f"✨ [Rerank-Test] 완료! 최고 점수: {top_hit.get('score', 0):.4f}")
        logger.info(f"🏆 [Rerank-Test] 최고 스니펫 (doc_id: {top_hit.get('doc_id', 'unknown')}): {top_hit.get('snippet', '')[:100]}...")

    context = "\n---\n".join(h["snippet"] for h in hits_sorted if h.get("snippet"))
    prompt = f"사용자 질의: {req.query}\n:\n{context}\n\n위 내용을 바탕으로 응답을 생성해 주세요."
    elapsed = round(time.perf_counter() - t0, 4)

    # 세션 파일명 체크(문서명 리스트)
    check_files: List[str] = []
    try:
        for h in hits_sorted:
            did = h.get("doc_id")
            if did:
                check_files.append(f"{did}.pdf")
            else:
                p = Path(h.get("path", ""))
                if str(p):
                    check_files.append(p.with_suffix(".pdf").name)
    except Exception:
        pass

    return {
        "elapsed_sec": elapsed,
        "settings_used": {"model": model_key, "searchType": search_type},
        "hits": [
            {
                "score": float(h["score"]),
                "path": h["path"],
                "chunk_idx": int(h["chunk_idx"]),
                "task_type": h["task_type"],
                "security_level": int(h["security_level"]),
                "doc_id": h.get("doc_id"),
                "snippet": h["snippet"],
            }
            for h in hits_sorted
        ],
        "prompt": prompt,
        "check_file": sorted(list(set(check_files))),
        "sid": sid,
        "collection": coll,
    }
# --- add: delete_test_files_by_names ---
async def delete_test_files_by_names(sid: str, file_names: List[str], task_type: Optional[str] = None):
    meta = get_test_session(sid)
    if not meta:
        return {"deleted": 0, "requested": len(file_names), "error": "invalid sid"}

    client = _client()
    coll = meta.get("collection") or _session_collection_name(sid)
    if coll not in client.list_collections():
        return {"deleted": 0, "requested": len(file_names), "error": "collection not found"}

    # 검증
    task_filter = ""
    if task_type:
        if task_type not in TASK_TYPES:
            return {"deleted": 0, "requested": len(file_names), "error": f"invalid taskType: {task_type}"}
        task_filter = f" && task_type == '{task_type}'"

    deleted_total = 0
    per_file: dict[str, int] = {}

    for name in (file_names or []):
        stem = Path(name).stem
        try:
            base_id, _ = _parse_doc_version(stem)
        except Exception:
            base_id = stem
        try:
            filt = f"doc_id == '{base_id}'{task_filter}"
            client.delete(collection_name=coll, filter=filt)
            deleted_total += 1
            per_file[name] = per_file.get(name, 0) + 1
        except Exception:
            logger.exception("[test-delete] failed: %s", name)
            per_file[name] = per_file.get(name, 0)

    try:
        client.flush(coll)
    except Exception:
        pass
    try:
        client.release_collection(collection_name=coll)
    except Exception:
        pass
    try:
        client.load_collection(collection_name=coll)
    except Exception:
        pass

    return {"deleted": deleted_total, "requested": len(file_names), "taskType": task_type, "perFile": per_file, "sid": sid}
    
# === 새 API: 키워드 없이 레벨 오버라이드 후 인제스트 ===
class OverrideLevelsRequest(BaseModel):
    """
    업로드(or 기존) 파일들에 대해 작업유형별 레벨을 강제로 세팅하고 인제스트.
    - files: 대상 파일 이름/경로(비우면 META 전체 대상이지만, 본 엔드포인트에서는 업로드 파일만 전달)
    - level_for_tasks: {"qna":2,"summary":1,"doc_gen":3} (필수)
    - tasks: 작업유형 제한 (미지정 시 모든 TASK_TYPES)
    """
    files: Optional[List[str]] = None
    level_for_tasks: Dict[str, int]
    tasks: Optional[List[str]] = None


async def override_levels_and_ingest(req: OverrideLevelsRequest):
    if not META_JSON_PATH.exists():
        return {"error": "메타 JSON이 없습니다. 먼저 /v1/admin/vector/extract 를 수행하세요."}

    target_tasks = [t for t in (req.tasks or TASK_TYPES) if t in TASK_TYPES]
    if not target_tasks:
        return {"error": "유효한 작업유형이 없습니다. (허용: doc_gen|summary|qna)"}

    level_map = {t: int(max(1, lv)) for t, lv in (req.level_for_tasks or {}).items() if t in TASK_TYPES}
    if not level_map:
        return {"error": "적용할 보안레벨이 없습니다. level_for_tasks 를 지정하세요."}

    meta = json.loads(META_JSON_PATH.read_text(encoding="utf-8"))

    # 대상 파일 셋(메타키/파일명/스텀 모두 허용)
    def _to_keyset(files: List[str]) -> set:
        out = set()
        for f in files:
            p = Path(f)
            out.update({str(f), p.name, p.stem})
        return out

    all_keys = list(meta.keys())  # "securityLevelX/.../파일명.확장자"
    if req.files:
        ks = _to_keyset(req.files)
        targets = [k for k in all_keys if (k in ks or Path(k).name in ks or Path(k).stem in ks)]
    else:
        targets = all_keys

    if not targets:
        return {"updated": 0, "ingested": 0, "message": "대상 파일이 없습니다."}

    updated = 0
    for k in targets:
        entry = meta.get(k) or {}
        sec = entry.get("security_levels") or {}
        for t in target_tasks:
            if t in level_map:
                sec[t] = int(level_map[t])
        entry["security_levels"] = sec
        meta[k] = entry
        updated += 1

    META_JSON_PATH.write_text(json.dumps(meta, ensure_ascii=False, indent=2), encoding="utf-8")

    # ★ 업로드한(또는 지정한) 파일만 인제스트
    res = await ingest_embeddings(
        model_key=None,
        chunk_size=None,
        overlap=None,
        target_tasks=target_tasks,
        collection_name=COLLECTION_NAME,
        file_keys_filter=targets,
    )
    return {
        "message": "레벨 오버라이드 후 인제스트 완료",
        "collection": COLLECTION_NAME,
        "updated_meta_entries": updated,
        "inserted_chunks": int(res.get("inserted_chunks", 0)),
        "target_count": len(targets),
    }